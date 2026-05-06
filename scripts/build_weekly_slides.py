"""Build Biweekly Meeting slides — Apr 27, 2026. QFL vs QE-SAC focus."""

import copy, sys, os
sys.path.insert(0, "/root/power-system")

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
import numpy as np
from lxml import etree

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn

TEMPLATE = "artifacts/templatesldie/Biweekly Personal Meeting Template.pptx"
OUTPUT   = "artifacts/templatesldie/Biweekly_Meeting_Apr27_2026.pptx"
DATE     = "April 27, 2026"

# Brand colours (match template dark blue + accents)
BLUE  = RGBColor(0x1F, 0x4E, 0x79)
GREEN = RGBColor(0x37, 0x86, 0x44)
RED   = RGBColor(0xC0, 0x00, 0x00)
GRAY  = RGBColor(0x60, 0x60, 0x60)


# ---------------------------------------------------------------------------
# Low-level helpers
# ---------------------------------------------------------------------------

def _clear_placeholder(shape):
    """Remove every paragraph/run from a placeholder text-frame."""
    tf = shape.text_frame
    tf.word_wrap = True
    for i in range(len(tf.paragraphs) - 1, 0, -1):
        p = tf.paragraphs[i]._p
        p.getparent().remove(p)
    tf.paragraphs[0].clear()


def fill_placeholder(shape, lines, size_pt=16):
    """Write lines into a placeholder, one paragraph per line."""
    _clear_placeholder(shape)
    tf = shape.text_frame
    tf.word_wrap = True
    first = True
    for line in lines:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size_pt)


def add_textbox(slide, left, top, width, height, lines, size_pt=14, color=None, bold=False):
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    for line in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size_pt)
        run.font.bold = bold
        if color:
            run.font.color.rgb = RGBColor(*color)
    return txBox


def fix_name_box(shape):
    tf = shape.text_frame
    lines = ["Ing Muyleang", "PhD Student — Quantum Computing Lab",
             "Pukyong National University", "muyleanging@gmail.com"]
    for i, p in enumerate(tf.paragraphs):
        if i >= len(lines):
            break
        for r in p.runs:
            r.text = ""
        if p.runs:
            p.runs[0].text = lines[i]
        else:
            p.add_run().text = lines[i]


# ---------------------------------------------------------------------------
# Figures
# ---------------------------------------------------------------------------

def fig_vqc_circuit():
    """Draw the 8-qubit VQC circuit diagram."""
    fig, ax = plt.subplots(figsize=(9, 3.8))
    ax.set_xlim(0, 10); ax.set_ylim(-0.5, 8.5)
    ax.axis("off")
    ax.set_title("VQC Circuit — 8 Qubits, 2 Layers (identical to Lin et al. 2025)", fontsize=11, pad=8)

    for q in range(8):
        y = 7 - q
        ax.axhline(y, xmin=0.02, xmax=0.98, color="black", lw=0.8, zorder=1)
        ax.text(-0.15, y, f"q{q}", ha="right", va="center", fontsize=9)

        # RY gate (state prep)
        ax.add_patch(plt.Rectangle((0.3, y-0.28), 0.9, 0.56, fill=True,
                                   facecolor="#AED6F1", edgecolor="#2980B9", lw=1.2, zorder=2))
        ax.text(0.75, y, f"RY(z{q})", ha="center", va="center", fontsize=7.5, zorder=3)

        # Layer 1: CNOT + RX
        if q < 7:
            ax.plot([2.1, 2.1], [y, y-1], "k-", lw=1.2, zorder=2)
            ax.plot(2.1, y, "ko", ms=6, zorder=3)
            ax.plot(2.1, y-1, "o", ms=8, color="white", markeredgecolor="black", lw=1.2, zorder=3)
            ax.plot([2.05, 2.15], [y-1, y-1], "k-", lw=1.2, zorder=4)
            ax.plot([2.1, 2.1], [y-1.15, y-0.85], "k-", lw=1.2, zorder=4)

        ax.add_patch(plt.Rectangle((2.9, y-0.28), 1.0, 0.56, fill=True,
                                   facecolor="#A9DFBF", edgecolor="#27AE60", lw=1.2, zorder=2))
        ax.text(3.4, y, f"RX(ζ{q}⁽¹⁾)", ha="center", va="center", fontsize=7, zorder=3)

        # Layer 2: CNOT + RX
        if q < 7:
            ax.plot([5.2, 5.2], [y, y-1], "k-", lw=1.2, zorder=2)
            ax.plot(5.2, y, "ko", ms=6, zorder=3)
            ax.plot(5.2, y-1, "o", ms=8, color="white", markeredgecolor="black", lw=1.2, zorder=3)
            ax.plot([5.15, 5.25], [y-1, y-1], "k-", lw=1.2, zorder=4)
            ax.plot([5.2, 5.2], [y-1.15, y-0.85], "k-", lw=1.2, zorder=4)

        ax.add_patch(plt.Rectangle((6.0, y-0.28), 1.0, 0.56, fill=True,
                                   facecolor="#A9DFBF", edgecolor="#27AE60", lw=1.2, zorder=2))
        ax.text(6.5, y, f"RX(ζ{q}⁽²⁾)", ha="center", va="center", fontsize=7, zorder=3)

        # Measurement
        ax.add_patch(plt.Rectangle((7.7, y-0.28), 0.9, 0.56, fill=True,
                                   facecolor="#FAD7A0", edgecolor="#E67E22", lw=1.2, zorder=2))
        ax.text(8.15, y, "⟨Z⟩", ha="center", va="center", fontsize=9, zorder=3)

    # Labels
    ax.text(0.75, -0.25, "State Prep\n(angle encoding)", ha="center", fontsize=8, color="#2980B9")
    ax.text(2.5, -0.25, "Layer 1\n(CNOT)", ha="center", fontsize=8)
    ax.text(3.4, -0.25, "RX L1", ha="center", fontsize=8, color="#27AE60")
    ax.text(5.5, -0.25, "Layer 2\n(CNOT)", ha="center", fontsize=8)
    ax.text(6.5, -0.25, "RX L2", ha="center", fontsize=8, color="#27AE60")
    ax.text(8.15, -0.25, "Measure", ha="center", fontsize=8, color="#E67E22")

    legend = [mpatches.Patch(color="#AED6F1", label="RY: angle encoding"),
              mpatches.Patch(color="#A9DFBF", label=f"RX: trainable (16 params)"),
              mpatches.Patch(color="#FAD7A0", label="⟨Z⟩: output ∈ [−1,1]")]
    ax.legend(handles=legend, loc="upper right", fontsize=8, framealpha=0.9)

    plt.tight_layout()
    path = "/tmp/fig_vqc.png"
    plt.savefig(path, dpi=150, bbox_inches="tight")
    plt.close()
    return path


def fig_architecture():
    """QFL architecture pipeline vs QE-SAC."""
    fig, axes = plt.subplots(1, 2, figsize=(10, 3.4))

    # --- QE-SAC (Lin et al.) ---
    ax = axes[0]
    ax.set_xlim(0, 6); ax.set_ylim(0, 1); ax.axis("off")
    ax.set_title("QE-SAC — Lin et al. (2025)\nSingle Utility, RL-based", fontsize=10, fontweight="bold", color="#C00000")

    boxes_l = [("obs\n[obs_dim]", 0.05, "#ECF0F1"),
               ("MLP\nencoder", 1.2, "#FDEBD0"),
               ("VQC\n8 qubits", 2.4, "#FDEBD0"),
               ("SAC\nactor", 3.6, "#FADBD8"),
               ("action", 4.8, "#ECF0F1")]
    for label, x, col in boxes_l:
        ax.add_patch(plt.Rectangle((x, 0.25), 0.9, 0.50, facecolor=col, edgecolor="gray", lw=1.2))
        ax.text(x+0.45, 0.50, label, ha="center", va="center", fontsize=8.5, fontweight="bold")
    for x in [0.95, 2.10, 3.30, 4.50]:
        ax.annotate("", xy=(x+0.25, 0.50), xytext=(x, 0.50),
                    arrowprops=dict(arrowstyle="->", color="black", lw=1.2))
    ax.text(3.0, 0.05, "Full actor + 2 critics + replay buffer ≈ 107 KB", ha="center",
            fontsize=8, color="#C00000", style="italic")

    # --- QFL (proposed) ---
    ax2 = axes[1]
    ax2.set_xlim(0, 6); ax2.set_ylim(0, 1); ax2.axis("off")
    ax2.set_title("QFL — This Work\n3 Utilities Federated, SPSA-based", fontsize=10, fontweight="bold", color="#1A5276")

    boxes_r = [("obs\n[obs_dim]", 0.05, "#ECF0F1"),
               ("Local\nEncoder\n(private)", 1.0, "#FDEBD0"),
               ("Shared\nHead+VQC\n(federated)", 2.15, "#D5F5E3"),
               ("Action\nHeads\n(private)", 3.4, "#FDEBD0"),
               ("action", 4.6, "#ECF0F1")]
    for label, x, col in boxes_r:
        ax2.add_patch(plt.Rectangle((x, 0.22), 1.0, 0.56, facecolor=col, edgecolor="gray", lw=1.2))
        ax2.text(x+0.5, 0.50, label, ha="center", va="center", fontsize=8, fontweight="bold")
    for x in [1.05, 2.15, 3.40, 4.60]:
        ax2.annotate("", xy=(x+0.05, 0.50), xytext=(x-0.05, 0.50),
                     arrowprops=dict(arrowstyle="->", color="black", lw=1.2))

    # Fed region highlight
    ax2.add_patch(plt.Rectangle((2.10, 0.18), 1.35, 0.64, fill=False,
                                edgecolor="#27AE60", lw=2.5, linestyle="--"))
    ax2.text(2.78, 0.10, "FedAvg: 280 params = 1.1 KB", ha="center",
             fontsize=8.5, color="#27AE60", fontweight="bold")
    ax2.text(2.78, 0.03, "383× less than classical FL", ha="center",
             fontsize=8, color="#27AE60", style="italic")

    plt.tight_layout()
    path = "/tmp/fig_arch.png"
    plt.savefig(path, dpi=150, bbox_inches="tight")
    plt.close()
    return path


def fig_reward():
    """Early QFL reward results + comparison bar."""
    fig, (ax1, ax2) = plt.subplots(1, 2, figsize=(10, 3.6))

    # Left: learning curves (partial data seed 0 local_only)
    rounds = [10, 20, 30]
    ax1.plot(rounds, [-0.60, -0.60, -0.60], "o-", label="13-bus (A)", color="#E74C3C", lw=2)
    ax1.plot(rounds, [-0.16, -0.18, -0.21], "s-", label="34-bus (B)", color="#3498DB", lw=2)
    ax1.plot(rounds, [-0.00, -0.00, -0.00], "^-", label="123-bus (C)", color="#2ECC71", lw=2)
    ax1.axhline(0, color="gray", ls="--", lw=0.8)
    ax1.set_xlabel("FL Round", fontsize=11)
    ax1.set_ylabel("Mean Reward / Step", fontsize=11)
    ax1.set_title("QFL Early Results\n(local_only, seed 0, rounds 10–30)", fontsize=10)
    ax1.legend(fontsize=10)
    ax1.set_xticks(rounds)
    ax1.grid(True, alpha=0.3)

    # Right: communication cost comparison
    methods = ["Classical\nSAC-FL", "QFL\n(This Work)"]
    kb = [419.0, 1.1]
    colors = ["#E74C3C", "#2ECC71"]
    bars = ax2.bar(methods, kb, color=colors, edgecolor="white", width=0.5)
    ax2.set_ylabel("KB per client per round (log scale)", fontsize=10)
    ax2.set_yscale("log")
    ax2.set_title("Communication Cost\n(383× reduction)", fontsize=10)
    for bar, val in zip(bars, kb):
        ax2.text(bar.get_x() + bar.get_width()/2, val*1.3,
                 f"{val:.1f} KB", ha="center", va="bottom", fontsize=11, fontweight="bold")
    ax2.annotate("", xy=(1, 1.1), xytext=(0, 419),
                 arrowprops=dict(arrowstyle="<->", color="#1A5276", lw=2))
    ax2.text(0.5, 30, "383×\nreduction", ha="center", fontsize=10,
             color="#1A5276", fontweight="bold")
    ax2.grid(True, axis="y", alpha=0.3)

    plt.tight_layout()
    path = "/tmp/fig_results.png"
    plt.savefig(path, dpi=150, bbox_inches="tight")
    plt.close()
    return path


# ---------------------------------------------------------------------------
# Slide construction
# ---------------------------------------------------------------------------

def find_shape(slide, name):
    for s in slide.shapes:
        if s.name == name:
            return s
    return None


def build():
    prs = Presentation(TEMPLATE)
    slides = prs.slides

    # ── Slide 0: Title ──────────────────────────────────────────────────────
    s0 = slides[0]
    sub = find_shape(s0, "Subtitle (Spring 2022)")
    runs = sub.text_frame.paragraphs[0].runs
    runs[0].text = DATE
    if len(runs) > 1:
        runs[1].text = ""
    fix_name_box(find_shape(s0, "[NAME]…"))

    # ── Slide 1: Index ──────────────────────────────────────────────────────
    fill_placeholder(find_shape(slides[1], "Text Placeholder 2"), [
        "Quantum Federated Learning (QFL)",
        "",
        "  1.  Research Topic & Motivation",
        "  2.  QFL vs QE-SAC — Method Comparison",
        "  3.  VQC Circuit & SPSA Training",
        "  4.  Implementation & Architecture",
        "  5.  Early Experiment Results",
        "  6.  Next Action Plan",
    ], size_pt=17)

    # ── Slide 2: Section divider ─────────────────────────────────────────────
    t2 = find_shape(slides[2], "Title 2")
    if t2.text_frame.paragraphs[0].runs:
        t2.text_frame.paragraphs[0].runs[0].text = "Quantum Federated Learning (QFL)"
    else:
        t2.text_frame.paragraphs[0].add_run().text = "Quantum Federated Learning (QFL)"

    # ── Slide 3: Executive Summary ───────────────────────────────────────────
    fill_placeholder(find_shape(slides[3], "Text Placeholder 4"), [
        "Problem: Lin et al. (2025) QE-SAC controls Volt-VAR for ONE utility — no federation.",
        "",
        "This Work (QFL): Extends QE-SAC to 3 utilities federated via FedAvg.",
        "  • Replaces SAC (RL) with SPSA — gradient-free VQC optimization, no critic, no replay buffer",
        "  • Federates only the QuantumEncoder (280 params = 1.1 KB) — private topology encoders stay local",
        "  • 383× communication reduction vs classical Federated SAC-FL",
        "",
        "Implementation: Complete. QFLAgent + QFLTrainer across 13-bus / 34-bus / 123-bus topologies.",
        "Experiment: Running (3 seeds × 2 conditions). Early results confirm stable training.",
    ], size_pt=16)

    # ── Slide 4: New Progress — Method Comparison ───────────────────────────
    fill_placeholder(find_shape(slides[4], "Text Placeholder 4"), [
        "QFL vs QE-SAC (Lin et al. 2025) — key differences:",
        "",
        "  Training:    QE-SAC uses SAC (RL)          | QFL uses SPSA (no RL, no critic)",
        "  Scope:       QE-SAC = 1 utility             | QFL = 3 utilities federated",
        "  Federated:   QE-SAC = No                   | QFL = Yes (FedAvg on QuantumEncoder)",
        "  Comm cost:   QE-SAC = N/A                  | QFL = 280 params = 1.1 KB/round",
        "  vs classical FL: —                          | 383× smaller than SAC-FL actor",
        "  VQC role:    QE-SAC = actor policy          | QFL = shared quantum encoder",
    ], size_pt=15)

    arch_path = fig_architecture()
    slides[4].shapes.add_picture(arch_path, Inches(0.5), Inches(3.8), Inches(9.0), Inches(3.0))
    add_textbox(slides[4], 0.5, 6.85, 9.0, 0.4,
                ["Figure 1: QFL vs QE-SAC architecture (federated region = green dashed box)"],
                size_pt=10, color=(80, 80, 80))

    # ── Slide 4b: New Progress (cont) — VQC + SPSA ─────────────────────────
    layout = prs.slide_layouts[1]
    s4b = prs.slides.add_slide(layout)
    for ph in s4b.placeholders:
        if ph.placeholder_format.idx == 0:
            ph.text = "New Progress (cont)"
        elif ph.placeholder_format.idx == 1:
            fill_placeholder(ph, [
                "VQC Circuit — identical to Lin et al. 2025 (fair comparison):",
                "  • 8 qubits, 2 layers, 16 trainable RX params",
                "  • State prep: RY(zᵢ), z = Tanh(SharedHead(obs)) × π ∈ (−π, π)",
                "  • Entanglement: CNOT(i, i+1) nearest-neighbour per layer",
                "  • Output: ⟨Zᵢ⟩ expectation per qubit → 8-dim feature",
                "",
                "SPSA Optimizer (replaces SAC backprop through VQC):",
                "  • δ ~ Rademacher(±1),  perturb θ ± c·δ,  evaluate reward twice",
                "  • ∂L/∂θ ≈ (L₊ − L₋) / (2c·δ)  — no quantum gradient needed",
                "  • Update: θ ← θ − α·∂L/∂θ    (c=0.1, α=0.01)",
                "  • Avoids barren plateau — SPSA never backprops through the circuit",
            ], size_pt=15)

    vqc_path = fig_vqc_circuit()
    s4b.shapes.add_picture(vqc_path, Inches(0.4), Inches(4.1), Inches(9.2), Inches(2.9))
    add_textbox(s4b, 0.4, 7.05, 9.2, 0.4,
                ["Figure 2: QFL VQC circuit — RY angle encoding, CNOT entanglement, RX trainable, ⟨Z⟩ measurement"],
                size_pt=10, color=(80, 80, 80))

    # ── Slide 4c: New Progress (cont 2) — Results ───────────────────────────
    s4c = prs.slides.add_slide(layout)
    for ph in s4c.placeholders:
        if ph.placeholder_format.idx == 0:
            ph.text = "New Progress (cont 2)"
        elif ph.placeholder_format.idx == 1:
            fill_placeholder(ph, [
                "Experiment: 3 seeds × 2 conditions (local_only vs qfl) — 🔄 Ongoing",
                "  Topologies: 13-bus (A, obs=43)  |  34-bus (B, obs=113)  |  123-bus (C, obs=349)",
                "  Rounds: 50  |  Steps/round: 1,000/client  |  Seeds: 0, 1, 2",
                "",
                "Early results — seed 0, local_only, mean reward/step:",
                "  Round 10:  A = −0.60  |  B = −0.16  |  C = −0.00",
                "  Round 20:  A = −0.60  |  B = −0.18  |  C = −0.00",
                "  Round 30:  A = −0.60  |  B = −0.21  |  C = −0.00",
                "",
                "Interpretation: C (123-bus) converges immediately — fewer critical buses per observation.",
                "B still learning. A stable but negative — voltage violations present, SPSA exploring.",
                "QFL (federated) condition starts after local_only completes per seed.",
            ], size_pt=15)

    res_path = fig_reward()
    s4c.shapes.add_picture(res_path, Inches(0.3), Inches(4.2), Inches(9.4), Inches(2.9))
    add_textbox(s4c, 0.3, 7.1, 9.4, 0.4,
                ["Figure 3 (left): QFL learning curves — Figure 3 (right): Communication cost (383× reduction vs classical FL)"],
                size_pt=10, color=(80, 80, 80))

    # Template indices: 5=Response(skip/remove), 6=Discussion, 7=NextAction
    # ── Slide 6 (Discussion) ─────────────────────────────────────────────────
    fill_placeholder(find_shape(slides[6], "Text Placeholder 4"), [
        "Items for advisor discussion:",
        "",
        "  1. Comparison baseline — QFL vs QE-SAC (Lin et al. 2025)",
        "     QE-SAC = single-utility RL. QFL = multi-utility, no RL, federated.",
        "     Metrics to compare: mean reward, violation rate, convergence speed.",
        "",
        "  2. Method name — 'QFL' (Quantum Federated Learning)",
        "     No RL, no SAC. Focus is on the FL + quantum representation.",
        "     Approval needed before paper title is finalized.",
        "",
        "  3. SPSA justification",
        "     SPSA avoids backprop through VQC → no barren plateau risk.",
        "     Tradeoff: 2 env evaluations per step (vs 1 for RL policy gradient).",
        "     Is this framing acceptable for IEEE TSG reviewers?",
    ], size_pt=16)

    # ── Slide 7 (Next Action) ─────────────────────────────────────────────────
    fill_placeholder(find_shape(slides[7], "Text Placeholder 4"), [
        "2-Week Plan:",
        "",
        "  Week 1 (Apr 27 – May 3):",
        "    • Collect full QFL results (3 seeds × 2 conditions) once run finishes",
        "    • Statistical analysis: local_only vs qfl — effect size d, t-test",
        "    • Add QE-SAC (Lin et al.) numbers as single-utility comparison baseline",
        "    • Generate final communication cost figure (383× bar chart)",
        "",
        "  Week 2 (May 4 – May 10):",
        "    • Write paper Section 3: QFL Methodology (SPSA + VQC + FedAvg)",
        "    • Write paper Section 4: Experiments (QFL vs QE-SAC vs local_only)",
        "    • Prepare VQC circuit diagram for paper (Figure 2 in paper)",
        "    • Internal draft → send to advisor for review",
    ], size_pt=16)

    # ── Remove "Response to Previous Feedback" slide (template index 5) ────
    # After the 3 dynamic slides are appended, template slide 5 is still in
    # position 5. We identify it by title and drop its XML + relationship.
    def remove_slide(prs, idx):
        xml_slides = prs.slides._sldIdLst
        slide = prs.slides[idx]
        # Find rId for this slide in presentation relationships
        for rel in prs.part.rels.values():
            if rel.reltype.endswith("/slide") and rel._target == slide.part:
                rId = rel.rId
                prs.part.drop_rel(rId)
                break
        sp_el = xml_slides[idx]
        xml_slides.remove(sp_el)

    # Template slide 5 = "Response to Previous Feedback" (0-indexed)
    remove_slide(prs, 5)

    # ── Reorder: move two new cont slides (now at positions -2,-1) to after slide[4] ──
    xml_slides = prs.slides._sldIdLst
    all_ids = list(xml_slides)
    s4b_xml = all_ids[-2]
    s4c_xml = all_ids[-1]
    xml_slides.remove(s4b_xml)
    xml_slides.remove(s4c_xml)
    xml_slides.insert(5, s4c_xml)
    xml_slides.insert(5, s4b_xml)

    prs.save(OUTPUT)
    print(f"\nSaved: {OUTPUT}")
    print("\nSlide summary:")
    for i, slide in enumerate(prs.slides):
        title = ""
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.name in ("Title", "Title 1", "Title 2", "Q & A"):
                title = shape.text.strip()[:60]
                break
        print(f"  {i+1:2d}. {title}")


if __name__ == "__main__":
    os.chdir("/root/power-system")
    build()
