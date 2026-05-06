"""
Generate Biweekly Meeting slides for April 15, 2026
Based on Biweekly-Personal-Meeting-Template.pptx
"""

import copy, os, json
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import numpy as np
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

BASE     = "/root/power-system"
TEMPLATE = f"{BASE}/artifacts/templatesldie/Biweekly Personal Meeting Template.pptx"
OUTPUT   = f"{BASE}/artifacts/templatesldie/Biweekly_Meeting_Apr15_2026.pptx"

# ─────────────────────────────────────────────
# Figure generation helpers
# ─────────────────────────────────────────────

def make_reward_bar_chart():
    path = "/tmp/figure_reward_comparison.png"
    conditions = ['Local-only', 'Naive FL', 'Aligned FL\n[PROPOSED]']
    clients = ['A (13-bus)', 'B (34-bus)', 'C (123-bus)']
    means = [
        [-6.569, -8.075, -7.093],
        [-6.663, -8.346, -7.170],
        [-6.597, -7.750, -7.077],
    ]
    stds = [
        [0.089, 0.521, 0.130],
        [0.162, 0.658, 0.084],
        [0.146, 0.621, 0.078],
    ]
    x = np.arange(len(clients))
    width = 0.25
    colors = ['#7f8c8d', '#e74c3c', '#2ecc71']

    fig, ax = plt.subplots(figsize=(8, 4.5))
    for i, (cond, m, s, c) in enumerate(zip(conditions, means, stds, colors)):
        bars = ax.bar(x + i*width, m, width, label=cond, color=c,
                      yerr=s, capsize=4, alpha=0.88, error_kw={'elinewidth':1.5})
    ax.set_xticks(x + width)
    ax.set_xticklabels(clients, fontsize=11)
    ax.set_ylabel('Mean Episode Reward (↑ better)', fontsize=10)
    ax.set_title('FL Reward Comparison — n=5 seeds, final 50 rounds', fontsize=11, fontweight='bold')
    ax.legend(fontsize=9)
    ax.set_ylim(-9.5, -6.0)
    # Annotate best results
    ax.annotate('d=+1.74\np=0.0089✓', xy=(1+2*width, -7.75), fontsize=8,
                ha='center', va='bottom', color='#27ae60', fontweight='bold')
    ax.annotate('d=+1.24\np=0.025', xy=(2+2*width, -7.077), fontsize=8,
                ha='center', va='bottom', color='#27ae60', fontweight='bold')
    ax.grid(axis='y', alpha=0.3)
    plt.tight_layout()
    plt.savefig(path, dpi=150, bbox_inches='tight')
    plt.close()
    return path

def make_comm_cost_chart():
    path = "/tmp/figure_comm_cost.png"
    methods = ['Classical\nSAC-FL', 'Base Paper\n(Lin 2025)', 'Naive\nQE-SAC-FL', 'Aligned FL\n[PROPOSED]']
    params  = [113288, 4896, 7480, 280]
    kb      = [453.2, 19.6, 29.9, 1.1]
    colors  = ['#e74c3c', '#e67e22', '#f39c12', '#2ecc71']

    fig, (ax1, ax2) = plt.subplots(1, 2, figsize=(9, 4))

    bars = ax1.bar(methods, params, color=colors, alpha=0.88)
    ax1.set_ylabel('Federated Parameters', fontsize=10)
    ax1.set_title('Federated Parameters per Round', fontsize=10, fontweight='bold')
    ax1.set_yscale('log')
    for bar, p in zip(bars, params):
        ax1.text(bar.get_x()+bar.get_width()/2, bar.get_height()*1.15,
                 f'{p:,}', ha='center', va='bottom', fontsize=8, fontweight='bold')
    ax1.annotate('405× smaller\nthan classical', xy=(3, 280), xytext=(2.5, 5000),
                arrowprops=dict(arrowstyle='->', color='green'),
                fontsize=8, color='green', fontweight='bold')

    bars2 = ax2.bar(methods, kb, color=colors, alpha=0.88)
    ax2.set_ylabel('KB per client per round', fontsize=10)
    ax2.set_title('Communication Cost (KB/round)', fontsize=10, fontweight='bold')
    ax2.set_yscale('log')
    for bar, k in zip(bars2, kb):
        ax2.text(bar.get_x()+bar.get_width()/2, bar.get_height()*1.15,
                 f'{k} KB', ha='center', va='bottom', fontsize=8, fontweight='bold')

    fig.suptitle('Communication Efficiency — Federated Setting', fontsize=11, fontweight='bold')
    plt.tight_layout()
    plt.savefig(path, dpi=150, bbox_inches='tight')
    plt.close()
    return path

def make_ablation_chart():
    path = "/tmp/figure_circuit_ablation.png"
    circuits = ['Paper\n(linear CNOT)', 'Deep\n(4-layer)', 'No Entangle', 'Ring CNOT']
    rewards  = [-5.879, -5.956, -5.991, -6.023]
    stds     = [0.082, 0.137, 0.051, 0.125]
    params   = [16, 32, 16, 16]
    colors   = ['#2ecc71', '#e67e22', '#e74c3c', '#9b59b6']

    fig, ax = plt.subplots(figsize=(7, 4))
    bars = ax.bar(circuits, rewards, color=colors, alpha=0.88,
                  yerr=stds, capsize=5, error_kw={'elinewidth':1.5})
    ax.set_ylabel('Mean Episode Reward (↑ better)', fontsize=10)
    ax.set_title('VQC Circuit Ablation — Client B, n=2 seeds', fontsize=11, fontweight='bold')
    ax.set_ylim(-6.4, -5.6)
    for bar, r, p in zip(bars, rewards, params):
        ax.text(bar.get_x()+bar.get_width()/2, bar.get_height()-0.02,
                f'{r:.3f}\n({p}p)', ha='center', va='top', fontsize=8, fontweight='bold', color='white')
    ax.grid(axis='y', alpha=0.3)
    ax.annotate('Best: linear CNOT\n2-layer = 16 params', xy=(0, -5.879),
                xytext=(1.2, -5.72), arrowprops=dict(arrowstyle='->', color='green'),
                fontsize=8, color='#27ae60', fontweight='bold')
    plt.tight_layout()
    plt.savefig(path, dpi=150, bbox_inches='tight')
    plt.close()
    return path

def make_gnn_vs_mlp_chart():
    path = "/tmp/figure_gnn_vs_mlp.png"
    clients = ['A (13-bus)', 'B (34-bus)', 'C (123-bus)']
    mlp_local   = [-5.932, -3.348, -6.977]
    mlp_aligned = [-5.983, -3.348, -6.993]
    gnn_local   = [-5.978, -3.978, -6.971]
    gnn_aligned = [-5.934, -3.873, -6.980]

    x = np.arange(len(clients))
    w = 0.2
    fig, ax = plt.subplots(figsize=(8, 4.5))
    ax.bar(x - 1.5*w, mlp_local,   w, label='MLP Local',    color='#3498db', alpha=0.8)
    ax.bar(x - 0.5*w, mlp_aligned, w, label='MLP Aligned',  color='#2ecc71', alpha=0.8)
    ax.bar(x + 0.5*w, gnn_local,   w, label='GNN Local',    color='#e67e22', alpha=0.8)
    ax.bar(x + 1.5*w, gnn_aligned, w, label='GNN Aligned',  color='#9b59b6', alpha=0.8)
    ax.set_xticks(x); ax.set_xticklabels(clients, fontsize=11)
    ax.set_ylabel('Mean Episode Reward (↑ better)', fontsize=10)
    ax.set_title('MLP vs GNN Encoder — n=3 seeds', fontsize=11, fontweight='bold')
    ax.legend(fontsize=9)
    ax.grid(axis='y', alpha=0.3)
    plt.tight_layout()
    plt.savefig(path, dpi=150, bbox_inches='tight')
    plt.close()
    return path

# ─────────────────────────────────────────────
# PPTX helpers
# ─────────────────────────────────────────────

def set_text(shape, text, font_size=None, bold=None, align=None):
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    if font_size:
        run.font.size = Pt(font_size)
    if bold is not None:
        run.font.bold = bold
    if align:
        p.alignment = align

def set_bullets(shape, bullets, font_size=14):
    """bullets: list of (text, level) tuples"""
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    first = True
    for text, level in bullets:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.level = level
        run = p.add_run()
        run.text = text
        run.font.size = Pt(font_size)

def add_image_to_slide(slide, img_path, left, top, width, height, caption=None):
    from pptx.util import Inches, Emu
    slide.shapes.add_picture(img_path, left, top, width, height)
    if caption:
        from pptx.util import Pt
        txBox = slide.shapes.add_textbox(left, top+height+Pt(4), width, Pt(20))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = caption
        run.font.size = Pt(9)
        run.font.color.rgb = RGBColor(0x66,0x66,0x66)
        p.alignment = PP_ALIGN.CENTER

def add_cont_slide(prs, template_slide_idx, title_text, bullets, img_path=None, img_caption=None):
    """Duplicate a slide layout and fill with continuation content"""
    template_slide = prs.slides[template_slide_idx]
    # Use the same layout
    layout = template_slide.slide_layout
    new_slide = prs.slides.add_slide(layout)

    # Set title
    for ph in new_slide.placeholders:
        if ph.placeholder_format.idx == 0:
            set_text(ph, title_text, font_size=28, bold=True)
        elif ph.placeholder_format.idx == 1:
            if img_path:
                # Mixed: bullets left, image right
                set_bullets(ph, bullets, font_size=12)
            else:
                set_bullets(ph, bullets, font_size=13)

    if img_path:
        sl_w = prs.slide_width
        sl_h = prs.slide_height
        img_left  = int(sl_w * 0.52)
        img_top   = int(sl_h * 0.18)
        img_w     = int(sl_w * 0.45)
        img_h     = int(sl_h * 0.65)
        add_image_to_slide(new_slide, img_path, img_left, img_top, img_w, img_h, img_caption)
    return new_slide

# ─────────────────────────────────────────────
# Main generation
# ─────────────────────────────────────────────

def add_new_slide(prs, ref_idx, title_text, bullets, img_path=None, img_caption=None, font_size=12, full_img=False):
    """Add a new slide using layout from ref_idx slide."""
    layout = prs.slides[ref_idx].slide_layout
    s = prs.slides.add_slide(layout)
    sl_w = prs.slide_width
    sl_h = prs.slide_height
    for ph in s.placeholders:
        if ph.placeholder_format.idx == 0:
            set_text(ph, title_text, font_size=26, bold=True)
        elif ph.placeholder_format.idx == 1:
            set_bullets(ph, bullets, font_size=font_size)
    if img_path and full_img:
        add_image_to_slide(s, img_path, int(sl_w*0.03), int(sl_h*0.42),
                           int(sl_w*0.94), int(sl_h*0.52), img_caption)
    elif img_path:
        add_image_to_slide(s, img_path, int(sl_w*0.52), int(sl_h*0.17),
                           int(sl_w*0.46), int(sl_h*0.65), img_caption)
    return s


def main():
    # Generate all figures first
    print("Generating figures...")
    fig_reward   = make_reward_bar_chart()
    fig_comm     = make_comm_cost_chart()
    fig_ablation = make_ablation_chart()
    fig_gnn      = make_gnn_vs_mlp_chart()
    fig_gradient = f"{BASE}/artifacts/qe_sac_fl/figures/vqc_gradient_instability.png"
    fig_4grid    = f"{BASE}/artifacts/qe_sac_fl/figures/vqc_4grid_circuits.png"
    print("Figures done.")

    prs = Presentation(TEMPLATE)
    sl_w = prs.slide_width
    sl_h = prs.slide_height
    # Template has 9 slides: 0=Title, 1=Index, 2=Research divider,
    # 3=Exec Summary, 4=New Progress, 5=Response, 6=Discussion, 7=Next Action, 8=Q&A
    content_layout = prs.slides[3].slide_layout  # "제목 및 내용" layout

    # ── Slide 0: Title ──
    s = prs.slides[0]
    for ph in s.placeholders:
        if ph.placeholder_format.idx == 21:
            set_text(ph, "April 15, 2026", font_size=20)
    for shape in s.shapes:
        if shape.name.startswith('[NAME]'):
            tf = shape.text_frame; tf.clear()
            for i, (t, bold, sz) in enumerate([
                ("Ing Muyleang", True, 16),
                ("Ph.D. Researcher", False, 13),
                ("Pukyong National University — Quantum Computing Lab", False, 12),
                ("muyleanging@pukyong.ac.kr", False, 11),
            ]):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.alignment = PP_ALIGN.CENTER
                r = p.add_run(); r.text = t
                r.font.bold = bold; r.font.size = Pt(sz)

    # ── Slide 1: Index ──
    s = prs.slides[1]
    for ph in s.placeholders:
        if ph.placeholder_format.idx == 1:
            set_bullets(ph, [
                ("QE-SAC-FL: Federated Quantum RL for Volt-VAR Control", 0),
                ("", 0),
                ("Experiments completed this week:", 0),
                ("Main FL comparison — n=5 seeds, 500 rounds (DONE)", 1),
                ("Personalized FL — fine-tuned per client (DONE)", 1),
                ("OpenDSS real-physics validation (DONE)", 1),
                ("GNN vs MLP encoder comparison (DONE — new)", 1),
                ("VQC circuit ablation — 4 variants (DONE — new)", 1),
                ("Latent alignment evaluation (Ongoing)", 1),
            ], font_size=13)

    # ── Slide 2: Research divider ──
    s = prs.slides[2]
    for ph in s.placeholders:
        if ph.placeholder_format.idx == 0:
            set_text(ph, "QE-SAC-FL Research", font_size=36, bold=True)

    # ── Slide 3: Executive Summary ──
    s = prs.slides[3]
    for ph in s.placeholders:
        if ph.placeholder_format.idx == 1:
            set_bullets(ph, [
                ("[1] Aligned encoder solves heterogeneous FL — Client B: d=+1.74, p=0.0089 ✓", 0),
                ("Client C: d=+1.24, p=0.0252 (trend)", 1),
                ("", 0),
                ("[2] 405× communication reduction (hard architectural fact)", 0),
                ("280 params = 1.1 KB/round  vs  classical SAC-FL 453 KB/round", 1),
                ("17.5× fewer params than base paper (Lin 2025)", 1),
                ("", 0),
                ("[3] heterogeneous FL gradient instability prevented", 0),
                ("Naive FL: 2.5× higher seed variance vs aligned FL", 1),
                ("", 0),
                ("Extended: Personalized FL — Client B: −3.261 vs local −8.075 (+59.6%)", 0),
                ("", 0),
                ("NEW this week: GNN encoder + VQC circuit ablation completed", 0),
            ], font_size=12)

    # ── Slide 4: New Progress (main results + reward chart) ──
    s = prs.slides[4]
    for ph in s.placeholders:
        if ph.placeholder_format.idx == 1:
            set_bullets(ph, [
                ("Main FL (n=5 seeds, 500 rounds):", 0),
                ("Aligned FL: A=−6.597, B=−7.750 ✓, C=−7.077 ✓", 1),
                ("Naive FL:   A=−6.663, B=−8.346,  C=−7.170", 1),
                ("Local-only: A=−6.569, B=−8.075,  C=−7.093", 1),
                ("", 0),
                ("GNN vs MLP (NEW — n=3 seeds):", 0),
                ("MLP aligned: A=−5.983, B=−3.348, C=−6.993", 1),
                ("GNN aligned: A=−5.934, B=−3.873, C=−6.980", 1),
                ("→ MLP wins static; GNN for dynamic topology (next)", 1),
                ("", 0),
                ("Personalized FL: Client B −3.261 (+59.6%)", 0),
                ("Latent alignment — Ongoing", 0),
            ], font_size=11)
    add_image_to_slide(s, fig_reward,
                       int(sl_w*0.52), int(sl_h*0.17),
                       int(sl_w*0.46), int(sl_h*0.55),
                       "Figure 1: Reward comparison n=5 seeds")

    # ── New Progress (cont) — comm cost ──
    add_new_slide(prs, 4, "New Progress (cont)",
        [("Communication cost — 405× reduction vs classical SAC-FL", 0),
         ("17.5× vs base paper | 3.3 MB vs 1.36 GB total bandwidth", 1),
         ("Hard architectural fact — no statistics needed", 1),
         ("", 0),
         ("VQC circuit ablation (4 variants tested):", 0),
         ("Paper (linear CNOT, 16p): −5.879 ✓ BEST", 1),
         ("No entanglement (16p):     −5.991 → entanglement matters", 1),
         ("Deep 4-layer (32p):        −5.956 → more params ≠ better", 1),
         ("Ring CNOT (16p):           −6.023 → worst", 1),
        ], img_path=fig_comm, img_caption="Figure 2: Comm cost (log scale)", font_size=11)

    # ── New Progress (cont 2) — GNN chart + gradient ──
    add_new_slide(prs, 4, "New Progress (cont 2)",
        [("GNN vs MLP encoder — MLP wins on static topology", 0),
         ("GNN designed for dynamic/fault scenarios (Task 32 next)", 1),
         ("Both: 280 federated params — identical communication cost", 1),
         ("", 0),
         ("VQC gradient instability under naive FL:", 0),
         ("Naive FL seed variance: 6.41×10⁻⁷ (2.5× > aligned)", 1),
         ("Naive FL round instability: 2.79×10⁻⁸ (1.4× > aligned)", 1),
         ("Mechanism: heterogeneous FL = directional conflict, not magnitude decay", 1),
        ], img_path=fig_gnn, img_caption="Figure 3: GNN vs MLP comparison", font_size=11)

    # ── Slide 5: Response to Feedback ──
    s = prs.slides[5]
    for ph in s.placeholders:
        if ph.placeholder_format.idx == 1:
            set_bullets(ph, [
                ("'Explain barren plateau' →", 0),
                ("Reframed as heterogeneous FL-induced gradient instability (data-supported, more accurate)", 1),
                ("Naive FL: 2.5× seed variance, 1.4× round instability vs aligned FL", 1),
                ("VQC circuit diagram generated (PennyLane qml.draw_mpl)", 1),
                ("", 0),
                ("'KB reduction — make it visual' →", 0),
                ("Bar chart + table created: 405× vs classical, 17.5× vs base paper", 1),
                ("1.36 GB → 3.3 MB total bandwidth over 500 rounds", 1),
                ("", 0),
                ("'Multiple grid topologies' →", 0),
                ("GNN vs MLP experiment: 3 topologies × 2 encoders × 3 conditions — DONE", 1),
                ("13-bus, 34-bus, 123-bus clients throughout all experiments", 1),
            ], font_size=12)

    # ── Slide 6: Discussion ──
    s = prs.slides[6]
    for ph in s.placeholders:
        if ph.placeholder_format.idx == 1:
            set_bullets(ph, [
                ("Prof concern: '280 params too small — classical already small'", 0),
                ("", 0),
                ("Answer: 280 params = FEDERATED communication cost per round", 1),
                ("Classical SAC-FL sends 113,288 params (453 KB) every round", 1),
                ("Ours = 405× smaller. Comparison is federated, not local model size.", 1),
                ("VQC: 16 params in 2⁸=256 dimensional Hilbert space — not weak", 1),
                ("", 0),
                ("Question: GNN as 4th contribution or ablation study only?", 0),
                ("GNN shows FL benefit on B (+0.105); MLP wins on static", 1),
                ("Recommend: ablation — 3 clean contributions is stronger for IEEE TSG", 1),
                ("", 0),
                ("Latent alignment eval finishing tonight — will update Table VIII", 0),
            ], font_size=12)

    # ── Slide 7: Next Action ──
    s = prs.slides[7]
    for ph in s.placeholders:
        if ph.placeholder_format.idx == 1:
            set_bullets(ph, [
                ("Tonight — Complete latent alignment eval → fill final metrics", 0),
                ("This week — Write paper Sections 3–5 (Related Work, Method, Results)", 0),
                ("This week — Generate all IEEE-format figures (300 DPI)", 0),
                ("Next week — Fault injection experiment (dynamic topology, GNN advantage)", 0),
                ("Next week — Advisor review of paper draft", 0),
                ("End of April — Submit to IEEE Transactions on Smart Grid", 0),
            ], font_size=13)

    # ── Slide 8: Q&A (leave as-is) ──

    # ── NEW: VQC 4-grid Circuit slide ──
    vqc_s = prs.slides.add_slide(content_layout)
    for ph in vqc_s.placeholders:
        if ph.placeholder_format.idx == 0:
            set_text(ph, "VQC Circuit Ablation — All 4 Variants", font_size=24, bold=True)
        elif ph.placeholder_format.idx == 1:
            set_bullets(ph, [
                ("Paper: Linear CNOT, 2-layer, 16p → −5.879 ✓ BEST  |  "
                 "No Ent: −5.991  |  Deep 4L: −5.956  |  Ring: −6.023", 0),
            ], font_size=11)
    add_image_to_slide(vqc_s, fig_4grid,
                       int(sl_w*0.02), int(sl_h*0.30),
                       int(sl_w*0.96), int(sl_h*0.64),
                       "Figure 4: All 4 VQC circuits (PennyLane draw_mpl) — 8 qubits each")

    # ── NEW: MLP vs GNN Architecture slide ──
    arch_s = prs.slides.add_slide(content_layout)
    for ph in arch_s.placeholders:
        if ph.placeholder_format.idx == 0:
            set_text(ph, "Architecture: MLP vs GNN LocalEncoder", font_size=24, bold=True)
        elif ph.placeholder_format.idx == 1:
            set_bullets(ph, [
                ("Both share: SharedEncoderHead [32→8, tanh×π] + VQC [8 qubits, 16p] = 280 federated params", 0),
                ("", 0),
                ("MLP LocalEncoder (original / this work baseline):", 0),
                ("obs [B, obs_dim]  →  Linear(obs_dim→64) → ReLU → Linear(64→32)  →  h [B,32]", 1),
                ("Flat features — no graph structure", 1),
                ("Aligned FL result: A=−5.983, B=−3.348, C=−6.993", 1),
                ("", 0),
                ("GNN LocalEncoder (new variant — topology-aware):", 0),
                ("obs → node features → BusGNN (graph conv × 2) → global mean pool → h [B,32]", 1),
                ("Uses power grid adjacency matrix — models bus connections directly", 1),
                ("Aligned FL result: A=−5.934, B=−3.873, C=−6.980", 1),
                ("", 0),
                ("Finding: MLP outperforms GNN on STATIC topology (matches Lee et al. 2022)", 0),
                ("GNN advantage expected in DYNAMIC topology / fault injection (Task 32 next)", 0),
            ], font_size=11)
    add_image_to_slide(arch_s, fig_gnn,
                       int(sl_w*0.52), int(sl_h*0.55),
                       int(sl_w*0.46), int(sl_h*0.38),
                       "Figure 5: MLP vs GNN aligned FL results")

    # ── Reorder: VQC + arch slides go before Q&A (slide 8) ──
    # Template: 9 slides (0-8). Added: 2 cont + vqc_grid + arch = 4 new = 13 total
    # 0=Title,1=Index,2=divider,3=ExecSum,4=NewProg,5=Response,6=Discussion,7=NextAction,8=Q&A
    # 9=cont1, 10=cont2, 11=vqc_grid, 12=arch
    # Target: 0,1,2,3,4,9,10,5,6,7,11,12,8
    sldIdLst = prs.slides._sldIdLst
    xml = list(sldIdLst)
    print(f"Total slides before reorder: {len(xml)}")
    target = [0,1,2,3,4,9,10,5,6,7,11,12,8]
    reordered = [xml[i] for i in target]
    for el in xml: sldIdLst.remove(el)
    for el in reordered: sldIdLst.append(el)

    # ─── Save ───
    prs.save(OUTPUT)
    print(f"\nSaved: {OUTPUT}")
    print(f"Slide count: {len(prs.slides)}")


if __name__ == "__main__":
    main()
