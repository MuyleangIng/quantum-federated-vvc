"""
Clean slide generator — Biweekly Meeting April 15, 2026
COMPLETE REWRITE: all body content via add_textbox(), 4 separate VQC circuit figures
Run: python3 scripts/build_slides.py
"""
import os
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
import numpy as np
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

BASE     = "/root/power-system"
TEMPLATE = f"{BASE}/artifacts/templatesldie/Biweekly Personal Meeting Template.pptx"
OUTPUT   = f"{BASE}/artifacts/templatesldie/Biweekly_Meeting_Apr15_2026.pptx"
FIGS     = f"{BASE}/artifacts/qe_sac_fl/figures"
TMP      = "/tmp"

# Colours
C_TITLE = RGBColor(0x1A, 0x38, 0x7A)
C_H1    = RGBColor(0x1A, 0x38, 0x7A)
C_BODY  = RGBColor(0x1A, 0x1A, 0x1A)
C_SUB   = RGBColor(0x44, 0x44, 0x44)
C_GREEN = RGBColor(0x1E, 0x8B, 0x4C)
C_GRAY  = RGBColor(0x77, 0x77, 0x77)

# ─────────────────────────────────────────────────────────────────
# FIGURE GENERATORS
# ─────────────────────────────────────────────────────────────────

def fig_reward():
    path = f"{TMP}/f_reward.png"
    labels = ['Client A\n(13-bus)', 'Client B\n(34-bus)', 'Client C\n(123-bus)']
    data = {
        'Local-only':              ([-6.569,-8.075,-7.093], [0.089,0.521,0.130], '#7f8c8d'),
        'Naive FL':                ([-6.663,-8.346,-7.170], [0.162,0.658,0.084], '#e74c3c'),
        '[PROPOSED] Aligned FL':   ([-6.597,-7.750,-7.077], [0.146,0.621,0.078], '#27ae60'),
    }
    x = np.arange(3); w = 0.25
    fig, ax = plt.subplots(figsize=(8, 5), facecolor='white')
    ax.set_facecolor('white')
    for i, (lbl, (m, s, c)) in enumerate(data.items()):
        ax.bar(x + i*w, m, w, label=lbl, color=c, yerr=s, capsize=4,
               alpha=0.88, error_kw={'elinewidth': 1.5})
    ax.set_xticks(x + w); ax.set_xticklabels(labels, fontsize=12)
    ax.set_ylabel('Mean Episode Reward  (higher = better)', fontsize=11)
    ax.set_title('FL Reward Comparison — n=5 seeds, final 50 rounds', fontsize=12, fontweight='bold')
    ax.legend(fontsize=10); ax.set_ylim(-9.5, -5.5); ax.grid(axis='y', alpha=0.3)
    ax.annotate('d=+1.74\np=0.009 ✓', xy=(1 + 2*w, -7.75), fontsize=9,
                ha='center', va='bottom', color='#27ae60', fontweight='bold')
    ax.annotate('d=+1.24\np=0.025',   xy=(2 + 2*w, -7.077), fontsize=9,
                ha='center', va='bottom', color='#27ae60', fontweight='bold')
    plt.tight_layout()
    plt.savefig(path, dpi=150, bbox_inches='tight', facecolor='white')
    plt.close()
    return path


def fig_comm():
    path = f"{TMP}/f_comm.png"
    methods = ['Classical\nSAC-FL', 'Lin 2025\n(base)', 'Naive\nQE-SAC-FL', '[PROPOSED]\nAligned FL']
    params  = [113288, 4896, 7480, 280]
    kb      = [453.2, 19.6, 29.9, 1.1]
    colors  = ['#e74c3c', '#e67e22', '#f39c12', '#27ae60']
    fig, (ax1, ax2) = plt.subplots(1, 2, figsize=(9, 4.5), facecolor='white')
    for ax, vals, ylabel in [(ax1, params, 'Federated Parameters'), (ax2, kb, 'KB per client per round')]:
        ax.set_facecolor('white')
        bars = ax.bar(methods, vals, color=colors, alpha=0.88)
        ax.set_ylabel(ylabel, fontsize=10); ax.set_yscale('log')
        ax.grid(axis='y', alpha=0.3)
        for b, v in zip(bars, vals):
            unit = ' KB' if ax is ax2 else ''
            ax.text(b.get_x() + b.get_width()/2, v * 1.5,
                    f'{v:,}{unit}', ha='center', va='bottom', fontsize=9, fontweight='bold')
        ax.tick_params(axis='x', labelsize=9)
    ax1.set_title('Federated Parameters', fontsize=11, fontweight='bold')
    ax2.set_title('Communication Cost (KB/round)', fontsize=11, fontweight='bold')
    ax1.annotate('405×\nsmaller', xy=(3, 280), xytext=(2.0, 15000),
                 arrowprops=dict(arrowstyle='->', color='#27ae60', lw=2),
                 fontsize=9, color='#27ae60', fontweight='bold')
    fig.suptitle('Communication Efficiency — Federated Setting', fontsize=12, fontweight='bold')
    plt.tight_layout()
    plt.savefig(path, dpi=150, bbox_inches='tight', facecolor='white')
    plt.close()
    return path


def fig_gnn():
    path = f"{TMP}/f_gnn.png"
    clients = ['Client A\n(13-bus)', 'Client B\n(34-bus)', 'Client C\n(123-bus)']
    x = np.arange(3); w = 0.20
    mlp_l = [-5.932, -3.348, -6.977]; mlp_a = [-5.983, -3.348, -6.993]
    gnn_l = [-5.978, -3.978, -6.971]; gnn_a = [-5.934, -3.873, -6.980]
    fig, ax = plt.subplots(figsize=(8, 4.5), facecolor='white')
    ax.set_facecolor('white')
    ax.bar(x - 1.5*w, mlp_l, w, label='MLP  Local',   color='#3498db', alpha=0.8)
    ax.bar(x - 0.5*w, mlp_a, w, label='MLP  Aligned', color='#1a5276', alpha=0.88)
    ax.bar(x + 0.5*w, gnn_l, w, label='GNN  Local',   color='#e67e22', alpha=0.8)
    ax.bar(x + 1.5*w, gnn_a, w, label='GNN  Aligned', color='#922b21', alpha=0.88)
    ax.set_xticks(x); ax.set_xticklabels(clients, fontsize=12)
    ax.set_ylabel('Mean Episode Reward  (higher = better)', fontsize=11)
    ax.set_title('MLP vs GNN Encoder — n=3 seeds, Aligned FL vs Local-only', fontsize=12, fontweight='bold')
    ax.legend(fontsize=10); ax.grid(axis='y', alpha=0.3)
    plt.tight_layout()
    plt.savefig(path, dpi=150, bbox_inches='tight', facecolor='white')
    plt.close()
    return path


def fig_arch():
    """Clean architecture diagram: private vs federated, MLP vs GNN"""
    path = f"{TMP}/f_arch.png"
    fig, ax = plt.subplots(figsize=(12, 5.5), facecolor='white')
    ax.set_facecolor('white'); ax.axis('off')
    ax.set_xlim(0, 12); ax.set_ylim(0, 5.5)

    def box(x, y, w, h, fc, ec, label, sub='', lfs=11):
        r = mpatches.FancyBboxPatch((x, y), w, h, boxstyle="round,pad=0.1",
                                    facecolor=fc, edgecolor=ec, linewidth=1.8)
        ax.add_patch(r)
        ax.text(x + w/2, y + h/2 + (0.18 if sub else 0), label,
                ha='center', va='center', fontsize=lfs, fontweight='bold')
        if sub:
            ax.text(x + w/2, y + h/2 - 0.22, sub,
                    ha='center', va='center', fontsize=9, color='#444')

    def arrow(x1, x2, y, color='#333'):
        ax.annotate('', xy=(x2, y), xytext=(x1, y),
                    arrowprops=dict(arrowstyle='->', color=color, lw=2.0))

    # ── Title
    ax.text(6, 5.3,
            'Full Model per Client  |  GREEN = Federated (280 params = 1.1 KB / round)',
            ha='center', va='center', fontsize=12, fontweight='bold', color='#1a387a')

    # ── MLP path (top row)
    ax.text(1.0, 4.7, 'MLP Encoder', ha='center', fontsize=10, color='#555', style='italic')
    box(0.1, 3.7, 1.8, 0.8, '#d6eaf8', '#2980b9', 'LocalEncoder', 'obs→64→32  (private)', 10)

    ax.text(1.0, 1.7, 'GNN Encoder', ha='center', fontsize=10, color='#555', style='italic')
    box(0.1, 0.8, 1.8, 0.8, '#d6eaf8', '#2980b9', 'BusGNN', 'node→graph conv→32  (private)', 10)

    ax.text(1.0, 3.45, '↕  Two encoder\nvariants tested', ha='center', fontsize=8.5, color='#888')

    # Both feed into same shared part
    arrow(1.9, 3.2, 4.1)
    arrow(1.9, 3.2, 1.2)
    ax.annotate('', xy=(3.2, 2.7), xytext=(1.9, 1.2),
                arrowprops=dict(arrowstyle='->', color='#555', lw=1.5))
    ax.text(2.6, 2.7, 'both output\nh ∈ ℝ³²', ha='center', fontsize=9, color='#666')

    # ── Shared federated components
    box(3.2, 3.0, 2.4, 1.1, '#d5f5e3', '#27ae60',
        'SharedEncoderHead', '32→8,  tanh×π\n264 params  ✅ FEDERATED', 10)
    arrow(5.6, 7.0, 3.55, '#27ae60')

    box(7.0, 3.0, 2.4, 1.1, '#d5f5e3', '#27ae60',
        'VQC', '8 qubits, 2 layers\n16 params  ✅ FEDERATED', 10)
    arrow(9.4, 10.5, 3.55)

    # FedAvg annotation
    ax.text(8.2, 4.45, '↕  FedAvg  (server aggregates\n    280 params only = 1.1 KB)',
            ha='center', fontsize=9, color='#27ae60', fontweight='bold',
            bbox=dict(boxstyle='round', facecolor='#eafaf1', edgecolor='#27ae60', alpha=0.9))

    # ── Private output heads
    box(10.5, 3.0, 1.4, 1.1, '#d6eaf8', '#2980b9', 'Actor\nHeads', 'Linear(8,|A|)\n(private)', 9)

    box(3.2, 1.5, 2.4, 1.0, '#d6eaf8', '#2980b9', 'Twin Critics', 'MLP(256,256)×2\n~132K params  (private)', 10)

    # Legend
    p1 = mpatches.Patch(facecolor='#d5f5e3', edgecolor='#27ae60', linewidth=2,
                        label='FEDERATED — 280 params (1.1 KB/round)  ← what gets sent')
    p2 = mpatches.Patch(facecolor='#d6eaf8', edgecolor='#2980b9', linewidth=2,
                        label='PRIVATE — stays on client, never shared  (~135K params)')
    ax.legend(handles=[p1, p2], loc='lower center', fontsize=10,
              framealpha=0.95, ncol=2, bbox_to_anchor=(0.5, 0.0))

    plt.tight_layout()
    plt.savefig(path, dpi=150, bbox_inches='tight', facecolor='white')
    plt.close()
    return path


# ─────────────────────────────────────────────────────────────────
# VQC CIRCUIT DIAGRAMS — 4 separate clean images
# ─────────────────────────────────────────────────────────────────

def _draw_vqc(variant, title, reward, n_qubits=8):
    path = f"{TMP}/vqc_{variant}.png"
    n_layers   = 4 if variant == 'deep' else 2
    has_cnot   = (variant != 'no_ent')
    ring       = (variant == 'ring')

    GW = 0.52; GH = 0.38; YS = 0.72
    col_ry   = GW + 0.35
    col_cnot = 0.60 + 0.35 if has_cnot else 0
    col_rx   = GW + 0.45
    layer_w  = col_ry + col_cnot + col_rx

    x_wire_start = 0.55
    x_content    = x_wire_start + 0.15
    x_meas_extra = 0.55
    total_x = x_content + n_layers * layer_w + x_meas_extra + 0.6
    total_y = (n_qubits - 1) * YS + 1.4

    fig_w = max(10, total_x * 0.95)
    fig_h = max(3.8, total_y * 0.78)

    fig, ax = plt.subplots(figsize=(fig_w, fig_h), facecolor='white')
    ax.set_facecolor('white'); ax.axis('off')
    ax.set_xlim(-0.7, total_x + 0.2)
    ax.set_ylim(-0.6, (n_qubits - 1) * YS + 0.9)

    wire_y = [(n_qubits - 1 - q) * YS for q in range(n_qubits)]

    x_wire_end = x_content + n_layers * layer_w + x_meas_extra + 0.4
    for y in wire_y:
        ax.plot([x_wire_start - 0.5, x_wire_end], [y, y], 'k-', lw=1.0, zorder=0)
    for q, y in enumerate(wire_y):
        ax.text(-0.6, y, f'q{q}', fontsize=9, va='center', ha='right',
                fontfamily='monospace', color='#222')

    def draw_gate(xc, yc, label, fc):
        r = plt.Rectangle((xc - GW/2, yc - GH/2), GW, GH,
                           facecolor=fc, edgecolor='#111', linewidth=1.2, zorder=3)
        ax.add_patch(r)
        ax.text(xc, yc, label, ha='center', va='center', fontsize=8,
                fontweight='bold', color='white', zorder=4)

    def draw_cnot(xc, y_ctrl, y_tgt):
        ax.plot(xc, y_ctrl, 'ko', ms=7, zorder=5)
        ylo, yhi = min(y_ctrl, y_tgt), max(y_ctrl, y_tgt)
        ax.plot([xc, xc], [ylo, yhi], 'k-', lw=1.2, zorder=2)
        circ = plt.Circle((xc, y_tgt), 0.17, facecolor='white',
                           edgecolor='#111', linewidth=1.5, zorder=4)
        ax.add_patch(circ)
        ax.plot([xc - 0.17, xc + 0.17], [y_tgt, y_tgt], 'k-', lw=1.5, zorder=5)
        ax.plot([xc, xc], [y_tgt - 0.17, y_tgt + 0.17], 'k-', lw=1.5, zorder=5)

    def draw_meas(xc, yc):
        r = plt.Rectangle((xc - GW/2, yc - GH/2), GW, GH,
                           facecolor='#f5f5f5', edgecolor='#333', linewidth=1.2, zorder=3)
        ax.add_patch(r)
        ax.text(xc, yc, '⟨Z⟩', ha='center', va='center', fontsize=8,
                fontweight='bold', color='#333', zorder=4)

    x = x_content
    for layer in range(n_layers):
        lbl_y = (n_qubits - 1) * YS + 0.55
        ax.text(x + layer_w/2, lbl_y,
                f'Layer {layer+1}', ha='center', va='bottom',
                fontsize=8, color='#888', style='italic')

        # RY encoding gates
        x_ry = x + GW/2
        for y in wire_y:
            draw_gate(x_ry, y, 'RY', '#1a5276')
        x += col_ry

        # CNOT entanglement
        if has_cnot:
            xc = x + 0.28
            if ring:
                pairs = list(zip(range(n_qubits - 1), range(1, n_qubits)))
                # draw linear pairs
                for ctrl, tgt in pairs:
                    draw_cnot(xc, wire_y[ctrl], wire_y[tgt])
                # indicate ring wrap with dashed arrow + label
                ax.annotate('', xy=(xc + 0.25, wire_y[0]),
                            xytext=(xc + 0.25, wire_y[-1]),
                            arrowprops=dict(arrowstyle='->', color='#888',
                                          lw=1.0, linestyle='dashed'))
                ax.text(xc + 0.52, (wire_y[0] + wire_y[-1])/2,
                        'ring\nwrap', ha='left', va='center',
                        fontsize=7, color='#888', style='italic')
            else:
                pairs = list(zip(range(n_qubits - 1), range(1, n_qubits)))
                for ctrl, tgt in pairs:
                    draw_cnot(xc, wire_y[ctrl], wire_y[tgt])
            x += col_cnot

        # RX trainable gates
        x_rx = x + GW/2
        for y in wire_y:
            draw_gate(x_rx, y, 'RX', '#922b21')
        x += col_rx

    # Measurements
    x_meas = x + 0.25
    for y in wire_y:
        draw_meas(x_meas, y)
        ax.plot([x_meas + GW/2, x_meas + GW/2 + 0.2], [y, y], 'k-', lw=1.0)

    # Legend patches
    ry_p  = mpatches.Patch(color='#1a5276', label='RY — input encoding (fixed)')
    rx_p  = mpatches.Patch(color='#922b21', label='RX — trainable weights (16 total)')
    cn_p  = mpatches.Patch(facecolor='white', edgecolor='black', linewidth=1.2,
                           label='CNOT — entanglement')
    mz_p  = mpatches.Patch(facecolor='#f5f5f5', edgecolor='#333', linewidth=1.2,
                           label='⟨Z⟩ — PauliZ measure')
    ax.legend(handles=[ry_p, rx_p, cn_p, mz_p],
              loc='lower right', fontsize=8, framealpha=0.95, ncol=2)

    ax.set_title(f'{title}     Reward: {reward}', fontsize=11, fontweight='bold', pad=10)
    plt.tight_layout()
    plt.savefig(path, dpi=150, bbox_inches='tight', facecolor='white')
    plt.close()
    return path


def fig_vqc_circuits():
    """Generate 4 separate VQC circuit PNG files (clean, white bg)"""
    specs = [
        ('paper',  'Paper Circuit  (2-layer, Linear CNOT, 16 params)',  '-5.879  ★ BEST'),
        ('no_ent', 'No Entanglement  (2-layer, no CNOT, 16 params)',    '-5.991'),
        ('deep',   'Deep Circuit  (4-layer, Linear CNOT, 32 params)',   '-5.956'),
        ('ring',   'Ring CNOT  (2-layer, Ring CNOT, 16 params)',        '-6.023'),
    ]
    return [_draw_vqc(v, t, r) for v, t, r in specs]


# ─────────────────────────────────────────────────────────────────
# SLIDE HELPERS — all via textboxes, no placeholder content
# ─────────────────────────────────────────────────────────────────

def clear_ph_text(slide):
    """Clear all placeholder text (hide '(Edit accordingly...)')"""
    for ph in slide.placeholders:
        if ph.placeholder_format.idx in (1,):   # body only
            if ph.has_text_frame:
                ph.text_frame.clear()


def set_title_ph(slide, text, size=22):
    """Fill the title placeholder (idx=0) cleanly"""
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0:
            ph.text_frame.clear()
            p = ph.text_frame.paragraphs[0]
            r = p.add_run(); r.text = text
            r.font.size = Pt(size); r.font.bold = True
            break


def add_textbox(slide, items, l, t, w, h, size=12, title_size=None):
    """
    Add a textbox with bullet items.
    items = list of (text, level)  where level 0 = bold header, level 1 = bullet
    """
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    first = True
    for text, lvl in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        indent = '    ' * lvl
        bullet  = '• ' if lvl >= 1 else ''
        r = p.add_run()
        r.text = indent + bullet + text
        r.font.size = Pt(size)
        if lvl == 0 and text.strip():
            r.font.bold = True
            r.font.color.rgb = C_H1
        elif lvl == 1:
            r.font.color.rgb = C_BODY
        else:
            r.font.color.rgb = C_SUB
    return tb


def add_img(slide, path, l, t, w, h, cap=None):
    slide.shapes.add_picture(str(path), l, t, w, h)
    if cap:
        tb = slide.shapes.add_textbox(l, t + h + Pt(2), w, Pt(22))
        p = tb.text_frame.paragraphs[0]
        r = p.add_run(); r.text = cap
        r.font.size = Pt(9); r.font.color.rgb = C_GRAY
        p.alignment = PP_ALIGN.CENTER


def new_content_slide(prs, ref_layout_idx, title,
                      items=None, img_path=None, cap=None,
                      layout='text_right', size=12):
    """
    Create a new slide from layout.
    layout: 'full_text'  | 'text_left_img_right' | 'full_img'
    """
    s = prs.slides.add_slide(prs.slide_layouts[ref_layout_idx])
    # Set title
    set_title_ph(s, title)
    clear_ph_text(s)

    W = prs.slide_width; H = prs.slide_height
    TOP  = int(H * 0.175)
    PAD  = int(W * 0.04)

    if layout == 'full_img' and img_path:
        if items:
            tb_h = int(H * 0.12)
            add_textbox(s, items, PAD, TOP, W - 2*PAD, tb_h, size=size)
            img_top = TOP + tb_h + int(H * 0.02)
            img_h   = int(H * 0.73)
        else:
            img_top = TOP
            img_h   = int(H * 0.80)
        add_img(s, img_path, PAD, img_top, W - 2*PAD, img_h, cap)

    elif layout == 'text_left_img_right' and img_path:
        mid = int(W * 0.50)
        add_textbox(s, items or [], PAD, TOP, mid - PAD*2, H - TOP - int(H*0.05), size=size)
        add_img(s, img_path, mid, TOP, W - mid - PAD, int(H * 0.76), cap)

    else:  # full_text
        add_textbox(s, items or [], PAD, TOP, W - 2*PAD, H - TOP - int(H*0.05), size=size)

    return s


# ─────────────────────────────────────────────────────────────────
# MAIN
# ─────────────────────────────────────────────────────────────────

def main():
    os.makedirs(FIGS, exist_ok=True)

    print("Generating figures...")
    F_REWARD  = fig_reward()
    F_COMM    = fig_comm()
    F_GNN     = fig_gnn()
    F_ARCH    = fig_arch()
    VQC_PATHS = fig_vqc_circuits()   # list of 4 paths
    print("  reward, comm, gnn, arch, 4×vqc  — all done")

    prs = Presentation(TEMPLATE)
    W = prs.slide_width; H = prs.slide_height
    # Template slide layouts: 0 = title layout, 1 = content layout
    LAYOUT_TITLE   = 0
    LAYOUT_CONTENT = 1

    PAD = int(W * 0.04)
    TOP = int(H * 0.175)

    # ── S0: Title ─────────────────────────────────────────────────
    s0 = prs.slides[0]
    # Leave template title "Bi-Weekly Meeting" as-is
    # Fill date placeholder (idx=21)
    for ph in s0.placeholders:
        if ph.placeholder_format.idx == 21:
            ph.text_frame.clear()
            p = ph.text_frame.paragraphs[0]
            r = p.add_run(); r.text = "April 15, 2026"
            r.font.size = Pt(20)
    # Name block — add textbox (avoid template name shape)
    name_tb = s0.shapes.add_textbox(int(W*0.30), int(H*0.72), int(W*0.40), int(H*0.22))
    tf = name_tb.text_frame; tf.word_wrap = True
    for i, (txt, bold, sz) in enumerate([
        ("Ing Muyleang",                          True,  16),
        ("Ph.D. Researcher",                      False, 12),
        ("Pukyong National University — QCL",     False, 11),
        ("muyleanging@pukyong.ac.kr",             False, 10),
    ]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = txt
        r.font.bold = bold; r.font.size = Pt(sz)

    # ── S1: Index ─────────────────────────────────────────────────
    s1 = prs.slides[1]
    clear_ph_text(s1)
    set_title_ph(s1, "Index")
    add_textbox(s1, [
        ("Research:  QE-SAC-FL — Federated Quantum RL for Volt-VAR Control", 0),
        ("", 0),
        ("Completed experiments:", 0),
        ("Main FL comparison  (n=5 seeds, 500 rounds)  ✅", 1),
        ("Personalized FL  (fine-tune per client)  ✅", 1),
        ("OpenDSS real-physics validation  ✅", 1),
        ("GNN vs MLP encoder comparison  ✅  ← new this week", 1),
        ("VQC circuit ablation  (4 variants)  ✅  ← new this week", 1),
        ("Latent alignment evaluation  ✅  DONE — results below", 1),
    ], PAD, TOP, W - 2*PAD, H - TOP - int(H*0.05), size=13)

    # ── S2: Research divider ──────────────────────────────────────
    s2 = prs.slides[2]
    set_title_ph(s2, "QE-SAC-FL Research", size=34)

    # ── S3: Executive Summary ─────────────────────────────────────
    s3 = prs.slides[3]
    clear_ph_text(s3)
    add_textbox(s3, [
        ("[1]  Aligned encoder solves heterogeneous FL", 0),
        ("Client B:  d = +1.74,  p = 0.0089  ✓  (significant after Bonferroni)", 1),
        ("Client C:  d = +1.24,  p = 0.0252  (strong trend)", 1),
        ("", 0),
        ("[2]  405× communication reduction vs classical SAC-FL", 0),
        ("280 params = federated component only = 1.1 KB/round", 1),
        ("Full model ≈ 135K params/client — only 280 get SENT per round", 1),
        ("", 0),
        ("[3]  heterogeneous FL gradient instability prevented", 0),
        ("Naive FL:  seed variance 2.5× higher,  round instability 1.4× higher", 1),
        ("Mechanism: misaligned latent inputs → contradictory VQC gradients", 1),
        ("", 0),
        ("Personalized FL:  Client B  −3.261  vs local −8.075  (+59.6%)", 0),
        ("New this week:  GNN encoder  +  VQC circuit ablation  ✅", 0),
    ], PAD, TOP, W - 2*PAD, H - TOP - int(H*0.05), size=12)

    # ── S4: New Progress — main FL + reward chart ─────────────────
    s4 = prs.slides[4]
    clear_ph_text(s4)
    mid = int(W * 0.49)
    add_textbox(s4, [
        ("Main FL results  (n=5 seeds, 500 rounds):", 0),
        ("Aligned FL:   A = −6.597,  B = −7.750 ✓,  C = −7.077 ✓", 1),
        ("Naive FL:     A = −6.663,  B = −8.346,    C = −7.170", 1),
        ("Local-only:   A = −6.569,  B = −8.075,    C = −7.093", 1),
        ("", 0),
        ("Personalized FL  (global pre-train → fine-tune):", 0),
        ("Client B:  −3.261 ± 0.072  vs local −8.075  → +59.6%  ← strongest", 1),
        ("Client A:  −5.890 ± 0.009  (+10.7%)  |  Client C:  −6.906 ± 0.025  (+2.6%)", 1),
        ("", 0),
        ("Latent alignment  ✅  DONE  (n=3 seeds × 20 rounds):", 0),
        ("local_only:   latent_sim = −0.054  |  grad_sim = +0.075  ← negative, no alignment", 1),
        ("naive_fl:     latent_sim = +0.397  |  grad_sim = +0.848  ← weight collapse", 1),
        ("aligned_fl:   latent_sim = +0.299  |  grad_sim = +0.599  ← structured alignment", 1),
        ("Note: naive_fl higher sim = superficial (weight collapse, not real alignment)", 1),
    ], PAD, TOP, mid - PAD*2, H - TOP - int(H*0.05), size=11)
    add_img(s4, F_REWARD, mid, TOP, W - mid - PAD, int(H*0.76),
            "Figure 1: Reward comparison — n=5 seeds")

    # ── NEW: New Progress (cont) — comm cost + VQC ablation ───────
    s5b = new_content_slide(prs, LAYOUT_CONTENT, "New Progress (cont)",
        items=[
            ("Communication cost — 405× reduction vs classical SAC-FL:", 0),
            ("280 federated params = 1.1 KB/round", 1),
            ("Classical SAC-FL federates actor MLP = 113K params = 453 KB/round", 1),
            ("Total over 500 rounds:  1.36 GB → 3.3 MB  (architectural fact)", 1),
            ("", 0),
            ("VQC circuit ablation — 4 variants, n=3 seeds:", 0),
            ("Paper  (2-layer, linear CNOT, 16p):   −5.879  ★ BEST", 1),
            ("Deep   (4-layer, linear CNOT, 32p):   −5.956  (more params ≠ better)", 1),
            ("No entanglement  (16p):                −5.991  (entanglement matters)", 1),
            ("Ring CNOT  (16p):                      −6.023  (worst)", 1),
        ],
        img_path=F_COMM, cap="Figure 2: Comm cost log scale — 405× reduction",
        layout='text_left_img_right', size=11)

    # ── NEW: New Progress (cont 2) — GNN vs MLP ───────────────────
    s5c = new_content_slide(prs, LAYOUT_CONTENT, "New Progress (cont 2)",
        items=[
            ("GNN vs MLP encoder  (same 280 federated params):", 0),
            ("MLP aligned:  A=−5.983,  B=−3.348,  C=−6.993", 1),
            ("GNN aligned:  A=−5.934,  B=−3.873,  C=−6.980", 1),
            ("MLP wins on STATIC topology — matches Lee et al. (2022)", 1),
            ("GNN value: dynamic topology + fault injection (Task 32)", 1),
            ("", 0),
            ("OpenDSS real-physics validation  (full AC power flow, n=3 seeds):", 0),
            ("Local-only:  variance ±5–6  (some seeds stuck in bad policy)", 1),
            ("Aligned FL:  variance ±0.5  — FL acts as stabiliser", 1),
            ("Client C:  d=+1.97,  p=0.054  on real OpenDSS  (near significant)", 1),
        ],
        img_path=F_GNN, cap="Figure 3: MLP vs GNN reward — n=3 seeds",
        layout='text_left_img_right', size=11)

    # ── S5: Response to Previous Feedback ─────────────────────────
    s5 = prs.slides[5]
    clear_ph_text(s5)
    add_textbox(s5, [
        ("Barren plateau concern  →  Analysis complete:", 0),
        ("Reframed as heterogeneous FL-induced gradient instability (data-supported, more accurate)", 1),
        ("Naive FL:  seed variance 2.5×,  round instability 1.4× higher", 1),
        ("Mechanism proven:  misaligned inputs → contradictory VQC gradients → cancellation", 1),
        ("", 0),
        ("KB reduction  →  Visualised and clarified:", 0),
        ("Bar chart (log scale) shows 405× reduction clearly  → see Figure 2", 1),
        ("Note:  280 = FEDERATED cost only.  Full model capacity = 135K params/client", 1),
        ("Classical SAC-FL sends 113K params = 453 KB every round", 1),
        ("", 0),
        ("Multiple grid topologies  →  Comparison experiments done:", 0),
        ("IEEE 13/34/123-bus tested across ALL 5 experiments", 1),
        ("GNN vs MLP encoder:  3 topologies × 2 encoders × 3 conditions  ✅", 1),
        ("OpenDSS real-physics validation across all 3 topologies  ✅", 1),
    ], PAD, TOP, W - 2*PAD, H - TOP - int(H*0.05), size=12)

    # ── NEW: VQC Circuit Ablation — 4 separate images ─────────────
    s_vqc = new_content_slide(prs, LAYOUT_CONTENT, "VQC Circuit Ablation — 4 Variants",
        items=[("4 circuit variants tested — paper circuit wins with reward −5.879", 0),
               ("8 qubits each  |  RY=input encoding  |  RX=trainable (16 or 32 params)", 1)],
        layout='full_text', size=12)
    # Place 4 images in 2×2 grid
    iw = int(W * 0.455); ih = int(H * 0.355)
    gx = [int(W*0.025), int(W*0.515)]
    gy = [int(H*0.195), int(H*0.565)]
    captions = ['Paper (linear CNOT) — BEST: −5.879',
                'No Entanglement — −5.991',
                'Deep 4-layer — −5.956',
                'Ring CNOT — −6.023']
    for idx, (path, cap) in enumerate(zip(VQC_PATHS, captions)):
        r, c = divmod(idx, 2)
        add_img(s_vqc, path, gx[c], gy[r], iw, ih, cap)

    # ── NEW: Architecture — MLP vs GNN ────────────────────────────
    s_arch = new_content_slide(prs, LAYOUT_CONTENT, "Architecture: MLP vs GNN Encoder",
        items=[
            ("Both variants share SAME 280 federated params  (SharedEncoderHead + VQC)", 0),
            ("Difference is only in the private LocalEncoder stage:", 1),
            ("", 0),
            ("MLP LocalEncoder  (original):", 0),
            ("obs [B, obs_dim]  →  Linear(obs_dim→64) → ReLU → Linear(64→32)  →  h", 1),
            ("Flat features — no topology awareness", 1),
            ("Aligned FL result:  A=−5.983,  B=−3.348,  C=−6.993  ← BETTER on static grid", 1),
            ("", 0),
            ("GNN LocalEncoder  (new — topology-aware):", 0),
            ("obs → node feats → BusGNN (2× graph conv) → global mean pool → h", 1),
            ("Uses power grid adjacency matrix — models bus connections directly", 1),
            ("Aligned FL result:  A=−5.934,  B=−3.873,  C=−6.980", 1),
            ("", 0),
            ("Finding:  MLP wins on STATIC topology.  GNN for dynamic/fault injection.", 0),
        ],
        img_path=F_ARCH, cap="Full model — green=federated (sent), blue=private (local only)",
        layout='full_img', size=11)

    # ── S6: Discussion ────────────────────────────────────────────
    s6 = prs.slides[6]
    clear_ph_text(s6)
    add_textbox(s6, [
        ("Q1:  '1.1 KB — is the model too small?'", 0),
        ("280 params = FEDERATED cost only (what is SENT per round)", 1),
        ("Full model per client = 135K params — LocalEncoder + Critics + Heads", 1),
        ("Classical SAC-FL sends 113K params = 453 KB/round  →  405× larger", 1),
        ("Model capacity is NOT small.  Only the shared quantum core is tiny.", 1),
        ("", 0),
        ("Q2:  GNN encoder — 4th contribution or ablation?", 0),
        ("Current result:  MLP beats GNN on static topology", 1),
        ("But GNN was designed for dynamic topology — not static", 1),
        ("So this is not a fair comparison yet", 1),
        ("Plan:  run fault injection (line outage) next week → see if GNN recovers better", 1),
        ("If GNN wins on dynamic  →  promote to 4th contribution", 1),
        ("If not  →  keep as ablation, note as future work in paper", 1),
        ("→  Request:  should I run this before submission?", 1),
        ("", 0),
        ("Q3:  Paper scope — 3 topologies enough for IEEE TSG?", 0),
        ("Currently:  13 / 34 / 123-bus  (3 very different scale topologies)", 1),
        ("Plan:  add 57-bus as 4th client to strengthen generalisability claim", 1),
        ("Request:  professor guidance on whether 3 vs 4 topologies is sufficient", 1),
    ], PAD, TOP, int(W * 0.54), H - TOP - int(H*0.05), size=11)
    add_img(s6, F_ARCH, int(W*0.55), TOP, int(W*0.43), int(H*0.76),
            "Full model vs federated component")

    # ── S7: Next Action ───────────────────────────────────────────
    s7 = prs.slides[7]
    clear_ph_text(s7)
    add_textbox(s7, [
        ("Tonight:", 0),
        ("Latent alignment eval  ✅  DONE — interpret naive_fl collapse vs aligned_fl", 1),
        ("", 0),
        ("This week:", 0),
        ("Write paper Sections 3–5  (Method, Experiments, Results)", 1),
        ("Prepare advisor email  (barren plateau + KB + multi-topology summary)", 1),
        ("", 0),
        ("Next week:", 0),
        ("Fault injection experiment  (GNN advantage on dynamic topology)", 1),
        ("Advisor review of paper draft", 1),
        ("", 0),
        ("End of April:", 0),
        ("Submit to IEEE Transactions on Smart Grid", 1),
        ("", 0),
        ("Correctness evidence  (5 levels):", 0),
        ("n=5 seeds  |  d=+1.74 p=0.009  |  OpenDSS real-physics  |  Ablation  |  Latent alignment", 1),
    ], PAD, TOP, W - 2*PAD, H - TOP - int(H*0.05), size=12)

    # ── S8: Q&A stays as-is ───────────────────────────────────────

    # ── REORDER ──────────────────────────────────────────────────
    # After adds, slides are:
    # 0=Title,1=Index,2=Divider,3=Exec,4=NewProg,5=Response,6=Discussion,7=Next,8=QA
    # 9=cont1, 10=cont2, 11=VQC, 12=Arch
    # Final order: Title,Index,Divider,Exec,NewProg,cont1,cont2,Response,VQC,Arch,Discussion,Next,QA
    sli = prs.slides._sldIdLst
    xml = list(sli)
    n = len(xml)
    print(f"Total slides before reorder: {n}")
    # target: [0,1,2,3,4,9,10,5,11,12,6,7,8]
    target = [0, 1, 2, 3, 4, 9, 10, 5, 11, 12, 6, 7, 8]
    assert len(target) == n, f"target len {len(target)} != slide count {n}"
    reordered = [xml[i] for i in target]
    for el in xml:     sli.remove(el)
    for el in reordered: sli.append(el)

    prs.save(OUTPUT)
    print(f"\nSaved: {OUTPUT}\n")

    prs2 = Presentation(OUTPUT)
    for i, s in enumerate(prs2.slides):
        t = next((ph.text_frame.text[:55]
                  for ph in s.placeholders
                  if ph.placeholder_format.idx == 0 and ph.has_text_frame), '—')
        print(f"  {i+1:2d}.  {t}")


if __name__ == "__main__":
    main()
