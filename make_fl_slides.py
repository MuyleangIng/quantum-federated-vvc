"""
Generate advisor meeting slides for QE-SAC-FL research.
Run:  python make_fl_slides.py
Out:  artifacts/QE_SAC_FL_AdvisorSlides.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

# ── Palette ────────────────────────────────────────────────────
BG      = RGBColor(0x0D, 0x1B, 0x2A)   # deep navy
BLUE    = RGBColor(0x1B, 0x6C, 0xA8)   # quantum blue
CYAN    = RGBColor(0x00, 0xC8, 0xFF)   # accent
GREEN   = RGBColor(0x2E, 0xCC, 0x71)   # pass / proven
ORANGE  = RGBColor(0xF3, 0x96, 0x1F)   # ongoing / warning
RED     = RGBColor(0xE7, 0x4C, 0x3C)   # fail / problem
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
GREY    = RGBColor(0xB0, 0xBE, 0xC5)
YELLOW  = RGBColor(0xF1, 0xC4, 0x0F)
PANEL   = RGBColor(0x14, 0x2A, 0x42)   # card background

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

os.makedirs("artifacts", exist_ok=True)

# ── Helpers ────────────────────────────────────────────────────
def bg(slide):
    r = slide.shapes.add_shape(1, 0, 0,
        prs.slide_width, prs.slide_height)
    r.fill.solid(); r.fill.fore_color.rgb = BG
    r.line.fill.background()

def rect(slide, l, t, w, h, color, line=None):
    s = slide.shapes.add_shape(1,
        Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = color
    if line:
        s.line.color.rgb = line
        s.line.width = Pt(1.5)
    else:
        s.line.fill.background()
    return s

def txt(slide, text, l, t, w, h,
        size=18, bold=False, color=WHITE,
        align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(
        Inches(l), Inches(t), Inches(w), Inches(h))
    tb.word_wrap = True
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    run = p.add_run(); run.text = text
    run.font.size = Pt(size); run.font.bold = bold
    run.font.color.rgb = color; run.font.italic = italic
    return tb

def pill(slide, label, l, t, color, size=13):
    rect(slide, l, t, 1.7, 0.32, color)
    txt(slide, label, l+0.08, t+0.03, 1.55, 0.28,
        size=size, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

def divider(slide, t):
    rect(slide, 0.4, t, 12.5, 0.04, CYAN)

# ══════════════════════════════════════════════════════════════
#  SLIDE 1 — TITLE
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); bg(s)
rect(s, 0, 0, 13.33, 0.12, CYAN)
rect(s, 0, 7.38, 13.33, 0.12, CYAN)

# glow box
rect(s, 1.5, 1.6, 10.3, 4.2, PANEL)
rect(s, 1.5, 1.6, 10.3, 4.2, PANEL, line=CYAN)

txt(s, "ADVISOR MEETING REPORT", 1.7, 1.8, 10, 0.6,
    size=13, bold=False, color=CYAN, align=PP_ALIGN.CENTER)
txt(s, "Federated Quantum RL for\nMulti-Utility Volt-VAR Control",
    1.7, 2.3, 10, 1.4, size=30, bold=True, color=WHITE,
    align=PP_ALIGN.CENTER)

divider(s, 3.85)

txt(s, "Three Novel Discoveries  ·  Experimental Proof  ·  Paper-Ready Results",
    1.7, 4.0, 10, 0.6, size=16, color=CYAN, align=PP_ALIGN.CENTER)

txt(s, "Ing Muyleang   |   Pukyong National University — QCL   |   April 1, 2026",
    1.7, 4.7, 10, 0.4, size=13, color=GREY, align=PP_ALIGN.CENTER)
txt(s, "Target: IEEE Transactions on Smart Grid  (IF 8.9)",
    1.7, 5.1, 10, 0.4, size=13, color=YELLOW, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════
#  SLIDE 2 — INDEX
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); bg(s)
rect(s, 0, 0, 13.33, 0.7, BLUE)
txt(s, "INDEX", 0.4, 0.1, 12, 0.55, size=24, bold=True,
    color=WHITE, align=PP_ALIGN.LEFT)

items = [
    ("01", "Executive Summary",   "What we built and the headline result",        CYAN),
    ("02", "New Progress",        "Three novel discoveries + experimental proof", GREEN),
    ("03", "Architecture",        "What gets federated vs what stays private",    BLUE),
    ("04", "All Results",         "Complete reward table — every condition",      CYAN),
    ("05", "Ongoing Work",        "What is running now",                          ORANGE),
    ("06", "Discussion",          "Why each finding is novel vs literature",      YELLOW),
    ("07", "Next Actions",        "Immediate tasks (2–3 weeks)",                  GREEN),
    ("08", "Plan",                "Full timeline to paper submission",            BLUE),
]

for i, (num, title, desc, color) in enumerate(items):
    row = i // 2; col = i % 2
    lx = 0.4 + col * 6.5; ty = 1.0 + row * 1.45
    rect(s, lx, ty, 6.1, 1.2, PANEL)
    rect(s, lx, ty, 0.55, 1.2, color)
    txt(s, num, lx+0.08, ty+0.32, 0.4, 0.5,
        size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(s, title, lx+0.65, ty+0.1, 5.2, 0.45,
        size=16, bold=True, color=WHITE)
    txt(s, desc, lx+0.65, ty+0.58, 5.2, 0.55,
        size=12, color=GREY)

# ══════════════════════════════════════════════════════════════
#  SLIDE 3 — EXECUTIVE SUMMARY
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); bg(s)
rect(s, 0, 0, 13.33, 0.7, BLUE)
txt(s, "01  EXECUTIVE SUMMARY", 0.4, 0.1, 12, 0.55,
    size=22, bold=True, color=WHITE)

# Problem box
rect(s, 0.4, 0.85, 5.9, 1.6, PANEL)
rect(s, 0.4, 0.85, 0.08, 1.6, RED)
txt(s, "THE PROBLEM", 0.6, 0.9, 5.5, 0.4, size=12, bold=True, color=RED)
txt(s, "3 utility companies need better Volt-VAR Control.\nThey CANNOT share raw grid data — competitors.",
    0.6, 1.28, 5.5, 1.0, size=14, color=WHITE)

# Solution box
rect(s, 6.9, 0.85, 5.9, 1.6, PANEL)
rect(s, 6.9, 0.85, 0.08, 1.6, GREEN)
txt(s, "THE SOLUTION", 7.1, 0.9, 5.5, 0.4, size=12, bold=True, color=GREEN)
txt(s, "Federated quantum RL — only 280 model parameters\ntravel to the server. Raw data never leaves the client.",
    7.1, 1.28, 5.5, 1.0, size=14, color=WHITE)

# 3 discoveries
rect(s, 0.4, 2.65, 3.8, 1.5, PANEL); rect(s, 0.4, 2.65, 0.08, 1.5, RED)
txt(s, "DISCOVERY 1 — QLSI", 0.6, 2.72, 3.5, 0.38, size=12, bold=True, color=RED)
txt(s, "Naive quantum FL makes\nevery client WORSE.\n(Not in any paper.)", 0.6, 3.1, 3.5, 0.95, size=13, color=WHITE)

rect(s, 4.7, 2.65, 3.8, 1.5, PANEL); rect(s, 4.7, 2.65, 0.08, 1.5, ORANGE)
txt(s, "DISCOVERY 2 — CSA", 4.9, 2.72, 3.5, 0.38, size=12, bold=True, color=ORANGE)
txt(s, "Pure aligned FL helps small\nclient early, large client late.\nNever all at once.", 4.9, 3.1, 3.5, 0.95, size=13, color=WHITE)

rect(s, 9.0, 2.65, 3.9, 1.5, PANEL); rect(s, 9.0, 2.65, 0.08, 1.5, YELLOW)
txt(s, "DISCOVERY 3 — PAD", 9.2, 2.72, 3.5, 0.38, size=12, bold=True, color=YELLOW)
txt(s, "Client dropout breaks\nalignment. Classical FL\nrobustness does not apply.", 9.2, 3.1, 3.5, 0.95, size=13, color=WHITE)

# Main result
rect(s, 0.4, 4.38, 12.5, 1.55, PANEL)
rect(s, 0.4, 4.38, 0.08, 1.55, GREEN)
txt(s, "MAIN RESULT — Personalised Federated Quantum RL", 0.6, 4.43, 12, 0.4,
    size=14, bold=True, color=GREEN)
txt(s, "13-bus  +50.2%     34-bus  +76.8%     123-bus  +24.8%     ALL THREE CLIENTS IMPROVE SIMULTANEOUSLY",
    0.6, 4.88, 12, 0.42, size=15, bold=True, color=WHITE)
txt(s, "Communication: 395× less data than federated classical SAC  ·  Mathematical proof  ·  No raw data shared",
    0.6, 5.33, 12, 0.42, size=13, color=CYAN)

# ══════════════════════════════════════════════════════════════
#  SLIDE 4 — NEW PROGRESS (discoveries)
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); bg(s)
rect(s, 0, 0, 13.33, 0.7, GREEN)
txt(s, "02  NEW PROGRESS — Three Novel Discoveries", 0.4, 0.1, 12, 0.55,
    size=22, bold=True, color=WHITE)

rows = [
    ("QLSI", "Quantum Latent Space Incompatibility",
     "Naive FedAvg on VQC only",
     "13-bus −336.6  34-bus −69.6  123-bus −5420.5\n(local: −331.4 / −65.5 / −5364.4)  ALL WORSE",
     "FedProx, SCAFFOLD address data heterogeneity — not encoder incompatibility",
     RED),
    ("CSA",  "Client Size Asymmetry",
     "Aligned FL 50 rounds vs 200 rounds",
     "Round 50: 13-bus PASS, 123-bus FAIL\nRound 200: 13-bus FAIL, 123-bus PASS  — benefit ROTATES",
     "Zhao 2018, Li 2021 study non-IID data content — not obs_dim gradient scale",
     ORANGE),
    ("PAD",  "Partial Alignment Drift",
     "2/3 clients per round, random dropout",
     "13-bus −341.4  34-bus −79.8  123-bus −5402.9\nAll worse AND VQC grad norms highest (0.000721) — oscillation",
     "McMahan 2017 proves dropout robustness for standard weights — not split-encoder coupling",
     YELLOW),
]

for i, (code, name, cond, result, gap, color) in enumerate(rows):
    y = 0.88 + i * 2.05
    rect(s, 0.3, y, 12.7, 1.85, PANEL)
    rect(s, 0.3, y, 0.5, 1.85, color)
    txt(s, code, 0.33, y+0.55, 0.45, 0.7,
        size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(s, name, 1.0, y+0.07, 5.0, 0.45, size=14, bold=True, color=color)
    txt(s, f"Condition: {cond}", 1.0, y+0.5, 4.8, 0.38, size=11, color=GREY, italic=True)
    txt(s, result, 1.0, y+0.88, 4.8, 0.75, size=12, color=WHITE)
    txt(s, "Gap in literature:", 6.2, y+0.07, 6.5, 0.35, size=11, bold=True, color=GREY)
    txt(s, gap, 6.2, y+0.45, 6.5, 1.2, size=12, color=CYAN, italic=True)

# ══════════════════════════════════════════════════════════════
#  SLIDE 5 — ARCHITECTURE
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); bg(s)
rect(s, 0, 0, 13.33, 0.7, BLUE)
txt(s, "03  ARCHITECTURE — Aligned QE-SAC-FL", 0.4, 0.1, 12, 0.55,
    size=22, bold=True, color=WHITE)

# Left: architecture flow
rect(s, 0.3, 0.85, 5.8, 6.3, PANEL)
txt(s, "ONE CLIENT (per utility)", 0.6, 0.95, 5.2, 0.4,
    size=13, bold=True, color=CYAN, align=PP_ALIGN.CENTER)

layers = [
    ("obs  (42 / 105 / 372 dim)", GREY, None, "INPUT"),
    ("LocalEncoder   obs→64→32", BLUE, "PRIVATE ❌  stays local", "4,832–25,952 params"),
    ("SharedEncoderHead   32→8", GREEN, "FEDERATED ✅  goes to server", "264 params"),
    ("8-qubit VQC   16 params", CYAN, "FEDERATED ✅  goes to server", "16 params"),
    ("Linear Head   8→132", BLUE, "LOCAL ❌  stays local", "1,188 params"),
    ("action  (0..131)", GREY, None, "OUTPUT"),
]

for i, (label, color, tag, params) in enumerate(layers):
    y = 1.5 + i * 0.85
    rect(s, 0.6, y, 5.2, 0.65, color if tag else PANEL)
    txt(s, label, 0.75, y+0.08, 3.0, 0.48, size=13, bold=True, color=WHITE)
    if tag:
        c = RED if "PRIVATE" in tag else GREEN
        txt(s, tag, 3.8, y+0.08, 1.8, 0.48, size=10, bold=True, color=c)
    txt(s, params, 0.75, y+0.42, 4.8, 0.25, size=9, color=YELLOW, italic=True)
    if i < len(layers)-1:
        txt(s, "▼", 3.1, y+0.65, 0.5, 0.22, size=11, color=GREY,
            align=PP_ALIGN.CENTER)

# Right: what travels vs stays
rect(s, 6.6, 0.85, 6.4, 3.1, PANEL)
txt(s, "WHAT TRAVELS TO SERVER", 6.8, 0.95, 6.0, 0.4,
    size=13, bold=True, color=GREEN)
travel = [
    ("SharedEncoderHead", "264 params", "1,056 bytes"),
    ("VQC weights",       " 16 params", "   64 bytes"),
    ("TOTAL per round",   "280 params", "1,120 bytes  ✅"),
]
for i, (name, p, b) in enumerate(travel):
    y = 1.45 + i*0.62
    txt(s, name, 6.8, y, 2.8, 0.55, size=13, color=WHITE)
    txt(s, p,    9.6, y, 1.5, 0.55, size=13, color=CYAN, align=PP_ALIGN.RIGHT)
    txt(s, b,   11.1, y, 1.7, 0.55, size=13, color=GREEN, align=PP_ALIGN.RIGHT)

rect(s, 6.6, 4.1, 6.4, 3.0, PANEL)
txt(s, "WHAT STAYS LOCAL (PRIVATE)", 6.8, 4.2, 6.0, 0.4,
    size=13, bold=True, color=RED)
stay = [
    ("LocalEncoder",  "feeder-specific compression"),
    ("Critics ×2",    "value estimation"),
    ("Replay buffer", "raw grid observations"),
    ("Linear head",   "action mapping"),
]
for i, (name, desc) in enumerate(stay):
    y = 4.72 + i*0.52
    txt(s, f"❌  {name}", 6.8, y, 3.2, 0.45, size=13, color=WHITE)
    txt(s, desc, 10.0, y, 3.0, 0.45, size=11, color=GREY, italic=True)

# Communication comparison
rect(s, 6.6, 7.1, 6.4, 0.6, BLUE)
txt(s, "Classical FL: 443,000 bytes/round   →   QE-SAC-FL: 1,120 bytes   =   395× LESS",
    6.7, 7.2, 6.2, 0.42, size=11, bold=True, color=YELLOW)

# ══════════════════════════════════════════════════════════════
#  SLIDE 6 — ALL RESULTS
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); bg(s)
rect(s, 0, 0, 13.33, 0.7, CYAN)
txt(s, "04  ALL RESULTS — Complete Reward Table", 0.4, 0.1, 12, 0.55,
    size=22, bold=True, color=BG)

headers = ["Condition", "13-bus", "34-bus", "123-bus", "All pass?", "Finding"]
widths  = [3.0, 1.5, 1.5, 1.8, 1.3, 3.4]
xs = [0.3]; [xs.append(xs[-1]+widths[i]) for i in range(len(widths)-1)]

rect(s, 0.3, 0.82, 12.7, 0.48, BLUE)
for i, h in enumerate(headers):
    txt(s, h, xs[i]+0.05, 0.87, widths[i]-0.1, 0.4,
        size=13, bold=True, color=WHITE)

rows = [
    ("Local only (baseline)",   "-331.4", "-65.5",  "-5364.4", "—",      "Reference",             PANEL, WHITE),
    ("Unaligned FL",            "-336.6", "-69.6",  "-5420.5", "NO ❌",  "← QLSI: all worse",     PANEL, RED),
    ("Aligned FL  50 rounds",   "-326.3", "-85.0",  "-5402.5", "NO ❌",  "← CSA: only 13-bus",    PANEL, ORANGE),
    ("Aligned FL  200 rounds",  "-339.5", "-69.3",  "-5251.4", "NO ❌",  "← CSA: only 123-bus",   PANEL, ORANGE),
    ("Partial FL  (2/3 clients)","-341.4","-79.8",  "-5402.9", "NO ❌",  "← PAD: all worse",      PANEL, YELLOW),
    ("Personalised FL  ★",      "-165.0", "-15.2",  "-4034.5", "YES ✅", "← SOLUTION",            GREEN, WHITE),
]

for i, (cond, a, b, c, p, note, bg_c, tc) in enumerate(rows):
    y = 1.38 + i * 0.82
    bg_color = RGBColor(0x14,0x2A,0x42) if bg_c==PANEL else GREEN
    rect(s, 0.3, y, 12.7, 0.75, bg_color)
    vals = [cond, a, b, c, p, note]
    for j, v in enumerate(vals):
        color = tc
        if j in [1,2,3] and p == "NO ❌": color = RED
        if j in [1,2,3] and p == "YES ✅": color = GREEN
        if j == 4 and p == "YES ✅": color = GREEN
        if j == 4 and p == "NO ❌": color = RED
        txt(s, v, xs[j]+0.05, y+0.15, widths[j]-0.1, 0.5, size=12,
            bold=(i==5), color=color)

# improvement row
rect(s, 0.3, 6.35, 12.7, 0.75, PANEL)
txt(s, "Improvement vs local:", 0.45, 6.45, 2.7, 0.55, size=13, bold=True, color=CYAN)
txt(s, "+50.2%", 3.35, 6.45, 1.4, 0.55, size=16, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
txt(s, "+76.8%", 4.85, 6.45, 1.4, 0.55, size=16, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
txt(s, "+24.8%", 6.35, 6.45, 1.7, 0.55, size=16, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
txt(s, "395× less communication (mathematical proof)", 8.1, 6.45, 4.8, 0.55,
    size=13, bold=True, color=YELLOW)

# ══════════════════════════════════════════════════════════════
#  SLIDE 7 — ONGOING WORK
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); bg(s)
rect(s, 0, 0, 13.33, 0.7, ORANGE)
txt(s, "05  ONGOING WORK", 0.4, 0.1, 12, 0.55,
    size=22, bold=True, color=WHITE)

ongoing = [
    ("Running now",
     "5-seed statistical significance runs\n"
     "All 6 conditions × 5 seeds → mean ± std error bars\n"
     "Overnight GPU job — 3× RTX 4090",
     GREEN, "HIGH"),
    ("Implementing",
     "Gradient-normalised FedAvg  (fix for CSA)\n"
     "Weight each client by 1/||∇S||  →  prevents large-feeder dominance\n"
     "Single-line change in fedavg_shared_head()",
     ORANGE, "HIGH"),
    ("Implementing",
     "FedProx regularisation on SharedHead  (fix for PAD)\n"
     "Add proximal term: (μ/2)||SharedHead − SharedHead_global||²\n"
     "Prevents SharedHead drift during partial-participation rounds",
     ORANGE, "HIGH"),
    ("Planned",
     "Transfer learning to 4th unseen feeder  (H7)\n"
     "Global VQC warm-start → new 4-bus feeder, zero prior training\n"
     "Connects directly to QE-SAC+ roadmap (Paper 2)",
     BLUE, "MEDIUM"),
]

for i, (status, desc, color, priority) in enumerate(ongoing):
    y = 0.88 + i * 1.52
    rect(s, 0.3, y, 12.7, 1.35, PANEL)
    rect(s, 0.3, y, 1.2, 1.35, color)
    txt(s, status, 0.33, y+0.22, 1.15, 0.9,
        size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(s, desc, 1.7, y+0.15, 9.5, 1.1, size=13, color=WHITE)
    pill(s, f"Priority: {priority}", 11.4, y+0.52, color, size=11)

# ══════════════════════════════════════════════════════════════
#  SLIDE 8 — DISCUSSION
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); bg(s)
rect(s, 0, 0, 13.33, 0.7, YELLOW)
txt(s, "06  DISCUSSION — Why Each Finding Is Novel", 0.4, 0.1, 12, 0.55,
    size=22, bold=True, color=BG)

# Header row
rect(s, 0.3, 0.82, 12.7, 0.45, BLUE)
for label, x, w in [("Finding", 0.35, 2.2), ("Closest Existing Paper", 2.6, 3.8),
                     ("What They Proved", 6.45, 3.2), ("What They MISSED", 9.7, 3.1)]:
    txt(s, label, x, 0.87, w-0.1, 0.38, size=12, bold=True, color=WHITE)

disc = [
    ("QLSI",
     "Li 2020 — FedProx\nKarimireddy 2020 — SCAFFOLD",
     "Heterogeneous data distributions\ncause client drift in FedAvg",
     "Encoder latent space incompatibility.\nHappens even with IID data.",
     RED),
    ("CSA",
     "Zhao 2018\nLi 2021 — FedBN",
     "Non-IID data content affects\nFedAvg convergence direction",
     "Gradient scale imbalance from\nobs_dim differences. Not data content.",
     ORANGE),
    ("PAD",
     "McMahan 2017 — FedAvg\nYang 2021 — partial participation",
     "FedAvg converges even when\nclients drop out each round",
     "Only proved for independent weights.\nNot for coupled split-encoder architecture.",
     YELLOW),
    ("H3 Proof",
     "No FL paper quantifies\ncommunication for quantum FL",
     "Classical FL requires full\nmodel weight exchange",
     "Quantum FL is 395–6920× cheaper.\nMathematical proof — not empirical.",
     GREEN),
]

for i, (finding, ref, proved, missed, color) in enumerate(disc):
    y = 1.38 + i * 1.42
    rect(s, 0.3, y, 12.7, 1.28, PANEL)
    rect(s, 0.3, y, 0.08, 1.28, color)
    txt(s, finding, 0.45, y+0.35, 2.0, 0.55, size=14, bold=True, color=color)
    txt(s, ref,     2.6,  y+0.1,  3.7, 1.0,  size=12, color=GREY, italic=True)
    txt(s, proved,  6.45, y+0.1,  3.1, 1.0,  size=12, color=WHITE)
    txt(s, missed,  9.7,  y+0.1,  3.1, 1.0,  size=12, color=CYAN)

# ══════════════════════════════════════════════════════════════
#  SLIDE 9 — NEXT ACTIONS
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); bg(s)
rect(s, 0, 0, 13.33, 0.7, GREEN)
txt(s, "07  NEXT ACTIONS  (2–3 Weeks)", 0.4, 0.1, 12, 0.55,
    size=22, bold=True, color=WHITE)

actions = [
    ("Week 1", "Statistical significance — 5 seeds",
     ["Run all 6 conditions × 5 seeds on 3 GPUs",
      "Compute mean ± std for every reward number",
      "Update all tables with error bars",
      "Expected: results confirmed, ±std < 10% of mean"],
     GREEN, "✅ Ready to run"),
    ("Week 1–2", "CSA fix — Gradient-normalised FedAvg",
     ["Modify fedavg_shared_head() in aligned_encoder.py",
      "Weight: w_i = 1 / max(||∇S_i||, 1e-8)",
      "Run 50r + 200r aligned FL with new weights",
      "Target: all 3 clients pass H1 simultaneously"],
     ORANGE, "🔄 Implementing"),
    ("Week 2", "PAD fix — FedProx on SharedHead",
     ["Add proximal term to aligned_agent.py actor loss",
      "L = L_SAC + (μ/2)||SharedHead − SharedHead_last||²",
      "Run partial FL (2/3 clients) with FedProx",
      "Target: partial FL passes H6 for all 3 clients"],
     ORANGE, "🔄 Implementing"),
    ("Week 3", "Math paragraph per finding for paper",
     ["QLSI: orthogonal latent bases argument (1 paragraph)",
      "CSA: gradient dominance inequality (1 paragraph)",
      "PAD: coupling constraint breaks McMahan proof (1 paragraph)",
      "Transforms observations into mechanisms → publishable"],
     BLUE, "📝 Writing"),
]

for i, (week, title, steps, color, status) in enumerate(actions):
    col = i % 2; row = i // 2
    x = 0.3 + col * 6.55; y = 0.88 + row * 3.15
    rect(s, x, y, 6.2, 2.95, PANEL)
    rect(s, x, y, 1.1, 0.5, color)
    txt(s, week, x+0.05, y+0.08, 1.0, 0.35,
        size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(s, title, x+1.25, y+0.08, 4.7, 0.38, size=13, bold=True, color=color)
    txt(s, status, x+1.25, y+0.48, 4.7, 0.32, size=11, color=GREY, italic=True)
    for j, step in enumerate(steps):
        txt(s, f"• {step}", x+0.15, y+0.85+j*0.5, 5.85, 0.45, size=11, color=WHITE)

# ══════════════════════════════════════════════════════════════
#  SLIDE 10 — PLAN / TIMELINE
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); bg(s)
rect(s, 0, 0, 13.33, 0.7, BLUE)
txt(s, "08  PLAN — Road to Submission", 0.4, 0.1, 12, 0.55,
    size=22, bold=True, color=WHITE)

phases = [
    ("APRIL 2026\nWeek 1–2",
     "Statistical runs + CSA fix + PAD fix",
     ["5-seed runs → mean ± std all conditions",
      "Gradient-normalised FedAvg → prove H1 all 3 clients",
      "FedProx on SharedHead → prove H6 with partial FL"],
     GREEN, "NOW"),
    ("APRIL 2026\nWeek 3–4",
     "Math justification + Paper draft",
     ["One math paragraph per novel finding (QLSI, CSA, PAD)",
      "Write Sections I–V: Intro, Background, Architecture, Results",
      "Generate publication-quality figures"],
     CYAN, "NEXT"),
    ("MAY 2026",
     "Transfer learning (H7) + QE-SAC+ bridge",
     ["4th unseen feeder transfer experiment",
      "Connect FL warm-start finding to QE-SAC+ roadmap",
      "Extend to GNN encoder if time permits"],
     ORANGE, "SOON"),
    ("JUNE 2026",
     "Paper revision + Submission",
     ["Internal review with advisor",
      "Address reviewer template checklist",
      "Submit: IEEE Transactions on Smart Grid (IF 8.9)"],
     BLUE, "TARGET"),
]

for i, (when, title, items, color, tag) in enumerate(phases):
    x = 0.25 + i * 3.2; y_top = 0.85
    rect(s, x, y_top, 3.0, 6.3, PANEL)
    rect(s, x, y_top, 3.0, 0.8, color)
    txt(s, when, x+0.1, y_top+0.05, 2.8, 0.65,
        size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(s, title, x+0.1, y_top+0.9, 2.8, 0.7,
        size=12, bold=True, color=color)
    for j, item in enumerate(items):
        txt(s, f"• {item}", x+0.1, y_top+1.65+j*0.85, 2.8, 0.75,
            size=11, color=WHITE)

    # connector arrow (not last)
    if i < 3:
        txt(s, "→", x+3.0, 3.7, 0.22, 0.45, size=22, color=GREY,
            align=PP_ALIGN.CENTER)

    # bottom tag
    rect(s, x, y_top+5.85, 3.0, 0.45, color)
    txt(s, tag, x+0.1, y_top+5.9, 2.8, 0.38,
        size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════
#  SLIDE 11 — CLOSING
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); bg(s)
rect(s, 0, 0, 13.33, 0.12, CYAN)
rect(s, 0, 7.38, 13.33, 0.12, CYAN)

rect(s, 1.5, 0.9, 10.3, 5.7, PANEL)
rect(s, 1.5, 0.9, 10.3, 5.7, PANEL, line=CYAN)

txt(s, "SIX CONTRIBUTIONS", 1.7, 1.0, 10, 0.45,
    size=14, bold=True, color=CYAN, align=PP_ALIGN.CENTER)
divider(s, 1.52)

contribs = [
    ("1", "QLSI",              "First identification: naive quantum FL hurts all clients", RED),
    ("2", "SharedEncoderHead", "Architecture fix: 280 federated params (395× less than classical)", GREEN),
    ("3", "CSA",               "First identification: gradient scale prevents simultaneous benefit", ORANGE),
    ("4", "PAD",               "First identification: partial dropout breaks encoder alignment", YELLOW),
    ("5", "Personalised QFL",  "Solution: +25–77% reward, all clients, 395× less communication", GREEN),
    ("6", "H3 Proof",          "Mathematical: 395–6920× communication reduction", CYAN),
]

for i, (num, name, desc, color) in enumerate(contribs):
    y = 1.65 + i * 0.72
    rect(s, 1.7, y, 0.4, 0.52, color)
    txt(s, num, 1.72, y+0.08, 0.36, 0.38,
        size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(s, name, 2.2, y+0.04, 3.0, 0.48, size=13, bold=True, color=color)
    txt(s, desc, 5.3, y+0.04, 6.2, 0.48, size=12, color=WHITE)

divider(s, 6.25)
txt(s, "IEEE Transactions on Smart Grid  ·  IF 8.9  ·  Submission: June 2026",
    1.7, 6.35, 10, 0.45, size=14, bold=True, color=YELLOW, align=PP_ALIGN.CENTER)

# ── Save ───────────────────────────────────────────────────────
out = "artifacts/QE_SAC_FL_AdvisorSlides.pptx"
prs.save(out)
print(f"Saved: {out}")
print(f"Slides: {len(prs.slides)}")
